#!/usr/bin/env python3
"""Generate multiple deep-analysis PPTX decks for ATD tech overview.

Usage:
    python3 docs/intro/build_decks.py
    # outputs: docs/intro/00-atd-overview.zh.pptx
    #         docs/intro/01-design-philosophy.zh.pptx
    #         docs/intro/02-architecture-deepdive.zh.pptx
    #         docs/intro/03-celia-phr-case-study.zh.pptx
    #         docs/intro/04-scenarios.zh.pptx

Theme: clean enterprise; Noto Sans CJK SC; 16:9 widescreen.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches, Emu

# ─────────────────────────── theme constants ────────────────────────────

FONT_CJK = "Noto Sans CJK SC"
FONT_MONO = "Noto Sans Mono CJK SC"

COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)       # dark blue
COLOR_ACCENT = RGBColor(0xE8, 0x5D, 0x04)        # orange
COLOR_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_MUTED = RGBColor(0x6B, 0x72, 0x80)
COLOR_BG_TINT = RGBColor(0xE0, 0xF2, 0xFE)       # pale blue
COLOR_BG_PANEL = RGBColor(0xF9, 0xFA, 0xFB)
COLOR_GOOD = RGBColor(0x16, 0xA3, 0x4A)
COLOR_BAD = RGBColor(0xDC, 0x26, 0x26)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.5)

# ──────────────────────────── helper layer ─────────────────────────────


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # entirely blank
    return prs.slides.add_slide(blank)


def _set_run(run, *, size=18, bold=False, italic=False,
             color=COLOR_TEXT, font=FONT_CJK):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, *, text="",
                size=18, bold=False, italic=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font=FONT_CJK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if text:
        run = p.add_run()
        run.text = text
        _set_run(run, size=size, bold=bold, italic=italic,
                 color=color, font=font)
    return tb, tf


def add_filled_rect(slide, left, top, width, height,
                    fill=COLOR_BG_PANEL, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    # remove default text in the shape
    shp.text_frame.text = ""
    return shp


def add_page_header(slide, *, deck_title, slide_title, page_num, total):
    # top thin band
    add_filled_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.65),
                    fill=COLOR_PRIMARY)
    # deck title (small, left)
    add_textbox(slide, Inches(0.5), Inches(0.12),
                Inches(6), Inches(0.4),
                text=deck_title, size=12, bold=True,
                color=COLOR_WHITE)
    # page number (right)
    add_textbox(slide, SLIDE_W - Inches(2.5), Inches(0.12),
                Inches(2), Inches(0.4),
                text=f"{page_num} / {total}",
                size=12, color=COLOR_WHITE,
                align=PP_ALIGN.RIGHT)
    # slide title
    add_textbox(slide, Inches(0.5), Inches(0.85),
                SLIDE_W - Inches(1.0), Inches(0.7),
                text=slide_title, size=28, bold=True,
                color=COLOR_PRIMARY)
    # divider line under title
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.55),
                                      SLIDE_W - Inches(0.5), Inches(1.55))
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(2)


def add_cover_slide(prs, *, deck_title, subtitle, footer):
    slide = add_blank_slide(prs)
    # full background panel
    add_filled_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H,
                    fill=COLOR_PRIMARY)
    # left accent bar
    add_filled_rect(slide, Emu(0), Emu(0), Inches(0.6), SLIDE_H,
                    fill=COLOR_ACCENT)
    # eyebrow
    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(0.6),
                text="ATD · Agent Tool Dispatch", size=18, bold=False,
                color=RGBColor(0xFB, 0xBF, 0x24))
    # title
    add_textbox(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(1.6),
                text=deck_title, size=46, bold=True,
                color=COLOR_WHITE)
    # subtitle
    add_textbox(slide, Inches(1.2), Inches(4.1), Inches(11), Inches(2.0),
                text=subtitle, size=22, color=RGBColor(0xCB, 0xD5, 0xE1),
                line_spacing=1.4)
    # footer
    add_textbox(slide, Inches(1.2), Inches(6.6), Inches(11), Inches(0.5),
                text=footer, size=13, italic=True,
                color=RGBColor(0xCB, 0xD5, 0xE1))
    return slide


def add_bullet_list(slide, left, top, width, height,
                    items: Iterable, *,
                    size=18, line_spacing=1.5,
                    bullet="•  ", color=COLOR_TEXT):
    """items: iterable of str OR tuples (text, sub_bullets:list[str])"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_top = Pt(2)
    first = True
    for item in items:
        sub = []
        if isinstance(item, tuple):
            text, sub = item[0], item[1]
        else:
            text = item
        # main bullet
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = bullet + text
        _set_run(run, size=size, color=color)
        # sub bullets
        for s in sub:
            sp = tf.add_paragraph()
            sp.alignment = PP_ALIGN.LEFT
            sp.line_spacing = line_spacing
            sp.level = 1
            r = sp.add_run()
            r.text = "    ▸  " + s
            _set_run(r, size=size - 3, color=COLOR_MUTED)
    return tb


def add_table(slide, left, top, width, height,
              header: list, rows: list,
              *, header_fill=COLOR_PRIMARY,
              header_color=COLOR_WHITE,
              cell_size=14,
              first_col_bold=True,
              col_widths=None):
    n_cols = len(header)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    # column widths
    if col_widths:
        total = sum(col_widths)
        for i, frac in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(width * frac / total))
    # header row
    for ci, h in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = str(h)
        _set_run(r, size=cell_size + 1, bold=True, color=header_color)
        cell.margin_left = Pt(6)
        cell.margin_right = Pt(6)
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
    # data rows
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                COLOR_BG_PANEL if ri % 2 == 1 else COLOR_WHITE
            )
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            _set_run(
                r, size=cell_size,
                bold=(first_col_bold and ci == 0),
                color=COLOR_TEXT,
            )
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return tbl_shape


def add_code_block(slide, left, top, width, height, code: str,
                   *, size=12):
    box = add_filled_rect(slide, left, top, width, height,
                          fill=RGBColor(0x11, 0x18, 0x27))
    tf = box.text_frame
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.word_wrap = True
    lines = code.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = line if line else " "
        _set_run(r, size=size, color=RGBColor(0xE5, 0xE7, 0xEB),
                 font=FONT_MONO)
    return box


def add_caption(slide, top, text, *, size=13, color=COLOR_MUTED,
                align=PP_ALIGN.LEFT, italic=True):
    add_textbox(slide, Inches(0.5), top,
                SLIDE_W - Inches(1.0), Inches(0.5),
                text=text, size=size, italic=italic, color=color,
                align=align)


def add_two_panel(slide, *, title_left, title_right):
    """Layout: two side-by-side panels under the title bar."""
    panel_top = Inches(1.8)
    panel_h = Inches(5.3)
    panel_w = (SLIDE_W - Inches(1.5)) / 2
    left_pos = Inches(0.5)
    right_pos = left_pos + panel_w + Inches(0.5)
    # backing rects
    add_filled_rect(slide, left_pos, panel_top, panel_w, panel_h,
                    fill=COLOR_BG_PANEL)
    add_filled_rect(slide, right_pos, panel_top, panel_w, panel_h,
                    fill=COLOR_BG_PANEL)
    # mini titles
    add_textbox(slide, left_pos + Inches(0.2), panel_top + Inches(0.15),
                panel_w - Inches(0.4), Inches(0.45),
                text=title_left, size=16, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, right_pos + Inches(0.2), panel_top + Inches(0.15),
                panel_w - Inches(0.4), Inches(0.45),
                text=title_right, size=16, bold=True, color=COLOR_PRIMARY)
    # body regions
    left_body = (left_pos + Inches(0.2), panel_top + Inches(0.75),
                 panel_w - Inches(0.4), panel_h - Inches(0.95))
    right_body = (right_pos + Inches(0.2), panel_top + Inches(0.75),
                  panel_w - Inches(0.4), panel_h - Inches(0.95))
    return left_body, right_body


def add_kv_grid(slide, left, top, width, height,
                pairs: list[tuple[str, str]],
                *, label_size=15, value_size=14):
    """Two-column key-value grid."""
    n = len(pairs)
    row_h = height / max(n, 1)
    for i, (k, v) in enumerate(pairs):
        y = top + Emu(int(row_h * i))
        # key
        add_textbox(slide, left, y,
                    Emu(int(width * 0.32)), Emu(int(row_h)),
                    text=k, size=label_size, bold=True,
                    color=COLOR_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
        # value
        add_textbox(slide, left + Emu(int(width * 0.32)), y,
                    Emu(int(width * 0.68)), Emu(int(row_h)),
                    text=v, size=value_size, color=COLOR_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE)


# ──────────────────────────── deck builders ─────────────────────────────


def build_deck_00_overview(out_path: Path):
    deck_title = "00 · ATD 总览"
    prs = new_presentation()

    add_cover_slide(
        prs,
        deck_title="ATD 总览 ── 是什么 / 为什么用 / 为什么不用 raw 替代品",
        subtitle=(
            "Agent Tool Dispatch — 跨 vendor 中立的 agent ↔ 工具调度协议\n"
            "1.1.0 已 publish · 17 个 crate · 4 个生产 adopter · Apache-2.0"
        ),
        footer="docs/intro/atd-tech-deck.zh.md · ATD maintainers · 2026-05",
    )

    # numbering: cover is 1
    slides_total = 17  # set after counting

    def hdr(s, title, n):
        add_page_header(s, deck_title=deck_title, slide_title=title,
                        page_num=n, total=slides_total)

    # 02 — 一页执行摘要
    s = add_blank_slide(prs)
    hdr(s, "一页执行摘要 — 不用 ATD vs 用 ATD", 2)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(5.2),
        header=["维度", "不用 ATD", "用 ATD"],
        rows=[
            ["工具被多 agent 平台调用",
             "每平台一套适配代码",
             "写一份 server,所有平台用"],
            ["多用户 / 多租户",
             "N 进程 × N 配置 × N OAuth",
             "一进程 + caller_id + 一份 broker"],
            ["审计 / 可观测性",
             "shell 历史 + grep stdout",
             "结构化 JSON Lines"],
            ["能力门禁",
             "各工具自检(漂移)",
             "dispatch 层一致 gate"],
            ["限流 / 超时",
             "各工具自实现",
             "tier deadline + semaphore"],
            ["跨厂商组合",
             "自写 multiplexer",
             "桥接多 socket,合并 catalog"],
            ["LLM 一次成功率",
             "看 --help(v1.1: 24%)",
             "structured(v1.2: 95.2%)"],
        ],
        cell_size=14,
        col_widths=[0.22, 0.36, 0.42],
    )
    add_caption(s, Inches(6.95),
                "实证 v1.4.0: ATD 路径 2 调用 ~1.6s 零错试 / CLI fallback 8 调用 ~6s 3 次走错",
                align=PP_ALIGN.LEFT)

    # 03 — 一句话定位
    s = add_blank_slide(prs)
    hdr(s, "一句话定位", 3)
    add_filled_rect(s, Inches(0.7), Inches(2.0),
                    SLIDE_W - Inches(1.4), Inches(3.5),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(1.0), Inches(2.3),
                SLIDE_W - Inches(2.0), Inches(2.9),
                text=(
                    "ATD 是 agent 调用工具时的一层中立调度协议。\n\n"
                    "Vendor 把工具 host 成一个 ATD server(Unix socket 或 HTTP),"
                    "任意 agent 平台(Hermes / Claude Code / Cursor / 自研)"
                    "通过同样的 wire 格式 discover / describe / call / dry-run。\n\n"
                    "中间层提供 capability gate / audit log / 多租户 token 路由 / "
                    "tool 可见性 / skill 同步 / cursor 分页 —— "
                    "raw CLI 拉不出、raw MCP 没规范、per-vendor 自研每个都要重写的东西。"
                ),
                size=20, color=COLOR_TEXT, line_spacing=1.5)
    add_caption(s, Inches(5.8),
                "完整定位见 docs/atd-positioning.md", align=PP_ALIGN.LEFT)

    # 04 — 四个任意
    s = add_blank_slide(prs)
    hdr(s, '"四个任意" —— 把工具世界的四种分裂折叠到一个统一面', 4)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.7),
        header=["维度", "现状的分裂", "ATD 的答案"],
        rows=[
            ["任意工具", "CLI / REST / MCP / native SDK 各自一套 shape",
             "一份 ToolDefinition 映射多个 binding"],
            ["任意平台", "Linux / macOS / iOS / Android / HMOS 各组调用面",
             "binding 选择在 server 侧 dispatch 时决定"],
            ["任意 agent", "Claude Code 吃不下 OpenAI function-calling shape",
             "所有 agent 同一份 SDK; adapter 渲 per-provider dict"],
            ["任意 framework", "LangChain ≠ MCP ≠ Apple App Intent",
             "一份定义,多 framework consumer"],
        ],
        cell_size=15,
        col_widths=[0.18, 0.42, 0.40],
    )
    add_caption(s, Inches(6.8),
                "源自 ATD architecture §1 的核心 interop claim",
                align=PP_ALIGN.LEFT)

    # 05 — v1.4 case study 数据
    s = add_blank_slide(prs)
    hdr(s, "v1.4.0 case study 实测 — ATD 严格优于 CLI fallback", 5)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.6),
        header=["维度", "ATD 路径", "CLI fallback 路径"],
        rows=[
            ["调用次数", "2", "8"],
            ["总耗时", "~1.6s", "~6s"],
            ["走错路径次数", "0", "3 (错 wrapper / --offset 不存在 ×2)"],
            ["首次拿到数据", "call #1 (1.2s)", "call #6 (5s)"],
            ["Audit 可观测性", "2 entries 完整", "shell log only"],
            ["agent 需自知 wrapper 命令", "否", "是 (healthkit healthkit +x 双关键字)"],
            ["agent 需自知 HMS 30 天上限", "否", "是 (撞错才知道)"],
        ],
        cell_size=14,
        col_widths=[0.30, 0.20, 0.50],
    )
    add_caption(s, Inches(6.7),
                "同一 Hermes session × DeepSeek × prompt; "
                "ATD bridge + CLI fallback 都摆在 agent 面前",
                align=PP_ALIGN.LEFT)

    # 06 — 三轮 case study progression
    s = add_blank_slide(prs)
    hdr(s, "healthkit_cli 三轮 case study —— v1.1 → v1.2 → v1.4", 6)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.0),
        header=["版本", "工具 surface", "LLM 表现"],
        rows=[
            ["v1.1.0",
             "8 个 raw HMS REST endpoint (permissive {type:object} schema)",
             "24% 成功率, 79 次调用, 66% Invalid param"],
            ["v1.2.0",
             "26 个 helper-tool (auto-derived 自 CLI + SKILL.md)",
             "95.2% 成功率, 21 次调用 (-73%), 1 次失败"],
            ["v1.4.0",
             "27 工具 + 多租户 mode",
             "2 ATD 调用 vs 8 CLI fallback, 0 错试"],
        ],
        cell_size=14,
        col_widths=[0.13, 0.45, 0.42],
    )
    add_caption(s, Inches(6.1),
                "schema + intent_examples + namespace = 成功率从 24% 到 95%",
                align=PP_ALIGN.LEFT)

    # 07 — 用 / 不用 ATD 的边界
    s = add_blank_slide(prs)
    hdr(s, "用 ATD / 不用 ATD 的边界", 7)
    left_body, right_body = add_two_panel(
        s, title_left="✅ 用 ATD 当", title_right="✘ 不用 ATD 当",
    )
    add_bullet_list(s, *left_body, items=[
        "工具 surface 要被 ≥1 个 LLM agent 平台用",
        "预期多 user / 多 caller_id (多租户)",
        "需要审计、可观测性、capability 门禁",
        "想给 Hermes + Claude Code + Cursor 用,不重复 N 套",
        "工具底下是真后端 (REST / DB / cloud API)",
        "多个 vendor 想一起 host 给同一 agent",
    ], size=15)
    add_bullet_list(s, *right_body, items=[
        "单进程脚本 + 单 user + 单工具",
        "工具是 sandbox 内纯计算 / 无 side effect",
        "工具在 agent 进程内 (in-process Tool trait 够了)",
        "你只要 MCP,且不需多租户/审计/跨 vendor",
    ], size=15)

    # 08 — vs raw CLI
    s = add_blank_slide(prs)
    hdr(s, "对比 raw CLI", 8)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.8),
        header=["关注点", "CLI", "ATD"],
        rows=[
            ["第一次调用成功率", "需先猜命令 / flag", "1 次成功, 0 retry"],
            ["多 agent 共享", "N 进程 / N 配置 / N OAuth",
             "1 server / N caller_id / broker 路由"],
            ["Audit", "shell history", "结构化 JSON Lines"],
            ["跨 vendor 组合", "自己写 multiplexer", "桥接多 socket"],
            ["LLM matching",
             "看 --help 文本 (混沌)",
             "description + intent_examples (结构化)"],
            ["升级安全",
             "flag 加减破坏 agent prompt",
             "tool def 是 schema, rev 跟踪"],
        ],
        cell_size=14,
        col_widths=[0.24, 0.38, 0.38],
    )

    # 09 — vs raw MCP
    s = add_blank_slide(prs)
    hdr(s, "对比 raw MCP — MCP 没规范的东西 ATD 在协议层 ship", 9)
    add_textbox(s, Inches(0.6), Inches(1.85),
                SLIDE_W - Inches(1.2), Inches(0.5),
                text="MCP 是 client-server 协议,缺以下能力 ——",
                size=15, italic=True, color=COLOR_MUTED)
    add_bullet_list(s, Inches(0.6), Inches(2.4),
                    SLIDE_W - Inches(1.2), Inches(4.3),
                    items=[
                        "server 侧 capability gate(每个 client 自己 gate)",
                        "server 侧 rate limit",
                        "multi-tenant token routing(MCP 假设单租户 stdio)",
                        "audit log 标准格式",
                        "tool visibility 多档(只有 hidden / visible 二元)",
                        "safety levels (Read / Write / Financial / Privacy / Physical / Destructive)",
                        "tier 概念 (Hot / Warm / Cold + 推导 deadline)",
                    ],
                    size=17)
    add_filled_rect(s, Inches(0.6), Inches(6.4),
                    SLIDE_W - Inches(1.2), Inches(0.85),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.8), Inches(6.5),
                SLIDE_W - Inches(1.6), Inches(0.65),
                text=(
                    "atd-mcp-bridge 兼容现有 MCP 客户端 — "
                    "Hermes / Claude Code / Cursor 不改一行代码即可接 ATD server。"
                ),
                size=14, italic=True, color=COLOR_PRIMARY,
                anchor=MSO_ANCHOR.MIDDLE)

    # 10 — vs 自研 adapter
    s = add_blank_slide(prs)
    hdr(s, "对比 per-vendor 自研 adapter", 10)
    add_textbox(s, Inches(0.6), Inches(1.95),
                SLIDE_W - Inches(1.2), Inches(2.5),
                text=(
                    "每写一次自研 adapter 都要重新设计:\n"
                    "  capability · audit · rate limit · token 管理 · stop logic\n\n"
                    "ATD 反过来 ——\n"
                    "  atd-runtime + atd-server ≈ 2000 行 Rust,\n"
                    "  vendor 写自己 server 只需:"
                ),
                size=17, color=COLOR_TEXT, line_spacing=1.5)
    add_code_block(s, Inches(0.6), Inches(4.6),
                   SLIDE_W - Inches(1.2), Inches(2.2),
                   code=(
                       "impl Tool for MyTool { fn definition() + fn call() }\n"
                       "Registry::register(my_tool);\n"
                       "atd_server::Server::new(registry, config).run().await?;\n\n"
                       "// healthkit_cli 的 `healthkit serve` 是 ~150 行 glue\n"
                       "// (一半是命令行参数解析)"
                   ),
                   size=14)

    # 11 — 5 message wire protocol
    s = add_blank_slide(prs)
    hdr(s, "Wire protocol — 5 message + 1 续传 (1.x 稳定面)", 11)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.6),
        header=["变体", "用途"],
        rows=[
            ["Hello", "握手。client_id / requested_capabilities / 可选 ucan_tokens"],
            ["Ping", "心跳"],
            ["ToolList",
             "Discovery。返回 Vec<ToolSummary>,按 DiscoverFilter 过滤"],
            ["ToolSchema", "单工具深 describe,返回完整 ToolDefinition"],
            ["RunTool",
             "调用。tool_id / args / CallOptions → ToolResultResponse"],
            ["RunToolContinue", "分页续传,带 opaque cursor"],
        ],
        cell_size=15,
        col_widths=[0.22, 0.78],
    )
    add_caption(s, Inches(6.7),
                "Schema 冻结为 1.x 稳定面 · UDS 和 HTTP 共用同 dispatch",
                align=PP_ALIGN.LEFT)

    # 12 — ToolDefinition 字段
    s = add_blank_slide(prs)
    hdr(s, "ToolDefinition —— 工具的完整声明", 12)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(5.4),
                   code=(
                       "pub struct ToolDefinition {\n"
                       "    pub id: String,                    // \"ref:fs.read\"\n"
                       "    pub name: String,\n"
                       "    pub description: String,           // LLM 看到的自然语言\n"
                       "    pub version: String,\n"
                       "    pub capability: ToolCapability,    // domain / actions / intent_examples\n"
                       "    pub input_schema: Value,           // JSON Schema 2020-12\n"
                       "    pub output_schema: Value,\n"
                       "    pub bindings: Vec<ToolBinding>,    // Native / Cli / future\n"
                       "    pub safety: ToolSafety,            // Read/Write/Financial/Privacy/...\n"
                       "    pub resources: ToolResources,      // timeout_ms / max_concurrent\n"
                       "    pub trust: ToolTrust,              // publisher / L0-L4\n"
                       "    pub visibility: ToolVisibility,    // Read / Write / Dangerous / Hidden\n"
                       "    pub required_capabilities: Vec<String>,\n"
                       "    pub tier: Option<ToolTier>,        // Hot / Warm / Cold\n"
                       "    pub errors: Vec<ToolErrorDef>,\n"
                       "}"
                   ),
                   size=14)

    # 13 — 5 层架构图
    s = add_blank_slide(prs)
    hdr(s, "5 层架构 + 3 核心 / 2 扩展", 13)
    layers = [
        ("User intent (voice · text · trigger)", COLOR_MUTED),
        ("Agent framework (Claude Code · Cursor · Hermes · LangChain)", COLOR_PRIMARY),
        ("Skills layer (adjacent) — SKILL.md · 渐进披露", COLOR_ACCENT),
        ("Client SDK — discover · describe · call · call_page · call_all", COLOR_PRIMARY),
        ("Dispatch — capability · tier · binding · cursor · middleware", COLOR_ACCENT),
        ("Listener — UDS (atd-server) / HTTP+MCP (atd-server-http)", COLOR_PRIMARY),
        ("Tool universe — bindings + 扩展点", COLOR_MUTED),
    ]
    top = Inches(1.95)
    h = Inches(0.55)
    gap = Inches(0.13)
    for i, (label, color) in enumerate(layers):
        y = top + Emu(int((h + gap) * i))
        bg = COLOR_PRIMARY if i in (1, 3, 5) else (
            COLOR_ACCENT if i in (2, 4) else RGBColor(0xCB, 0xD5, 0xE1))
        fg = COLOR_WHITE if i in (1, 2, 3, 4, 5) else COLOR_TEXT
        add_filled_rect(s, Inches(0.7), y, SLIDE_W - Inches(1.4), h, fill=bg)
        add_textbox(s, Inches(0.9), y, SLIDE_W - Inches(1.8), h,
                    text=label, size=16, bold=True, color=fg,
                    anchor=MSO_ANCHOR.MIDDLE)
    add_caption(s, Inches(6.7),
                "3 core = schema · dispatch · security; "
                "2 extension = bindings · middleware",
                align=PP_ALIGN.LEFT)

    # 14 — 17 crate map
    s = add_blank_slide(prs)
    hdr(s, "17 crate workspace map (1.1.0 stable)", 14)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(5.3),
        header=["crate", "职责"],
        rows=[
            ["atd-protocol", "wire 格式 + 类型 + sanitize"],
            ["atd-sdk", "Rust 客户端 SDK (discover/describe/call/call_page/call_all)"],
            ["atd-runtime",
             "registry / dispatch / audit / rate / TokenBroker / UCAN / Cursor / Metrics"],
            ["atd-server", "Unix socket listener + 连接任务"],
            ["atd-server-http",
             "HTTP listener + MCP JSON-RPC translator + bearer auth + SSE"],
            ["atd-middleware-fhir", "FHIR R4 egress validation (75-URI 白名单)"],
            ["atd-middleware-pii-redact-medical",
             "HIPAA Safe Harbor PHI redaction (18 类 × 13 路径)"],
            ["atd-tools-{echo, fs, shell, web}", "4 个内置工具示例 crate"],
            ["atd-mcp-bridge", "MCP/stdio ↔ ATD wire 桥"],
            ["atd-cli", "atd 开发者 CLI"],
            ["atd-ref-server", "参考 server binary"],
            ["atd-conformance", "跨实现 conformance fixture"],
            ["atd-mock-weather-server",
             "跨 vendor demo bin (publish = false)"],
        ],
        cell_size=12,
        col_widths=[0.30, 0.70],
    )

    # 15 — 适用场景一览
    s = add_blank_slide(prs)
    hdr(s, "7 大高价值应用场景", 15)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.3),
                    items=[
                        ("医疗 / PHR 合规",
                         ["FHIR + HIPAA PHI · UCAN sub-agent delegation · 多 OAuth 多用户"]),
                        ("Agent-Native CLI",
                         ["一个 binary 四个出口 · CLI/UDS/MCP-stdio/MCP-HTTP"]),
                        ("跨厂商工具组合 (cross-vendor)",
                         ["桥接多 socket · 每 vendor 自治 audit + broker"]),
                        ("Embodied agent / 物理仿真",
                         ["Python in-process server · shared world state 显式声明"]),
                        ("跨设备 federation (远程 ATD endpoint)",
                         ["cursor 分页 + provenance · CRDT 同步 + 失败恢复"]),
                        ("Agentic IDE / 代码助手",
                         ["统一 fs / shell / web 表面 · audit · capability 升权"]),
                        ("多 agent 编排 (orchestrator + N children)",
                         ["UCAN delegation chain · caller_id 路由 + audit"]),
                    ],
                    size=16, line_spacing=1.35)

    # 16 — Adopter 实例
    s = add_blank_slide(prs)
    hdr(s, "Active adopters (post-1.0)", 16)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.6),
        header=["Adopter", "Transport", "驱动的 SP", "状态"],
        rows=[
            ["healthkit_cli", "Unix socket",
             "首个 vendor server · 多租户 broker",
             "case study v1.4.0 头对头"],
            ["celia_phr", "HTTP via atd-server-http",
             "streamable-http · token-broker-phase2 · capability-v2 · "
             "medical-middleware · concurrency-baseline",
             "closed-verified"],
            ["cbrain", "Python in-process",
             "server-py-v1 · 22/24 conformance",
             "adopter-validation"],
            ["oh-cli / mycli", "Unix + HTTP",
             "mcp-native-v1 · CLI ↔ ATD bidirectional",
             "1.1.0 同日落地"],
        ],
        cell_size=12,
        col_widths=[0.16, 0.20, 0.38, 0.26],
    )
    add_caption(s, Inches(6.7),
                "ATD post-1.0 SP 触发模式: 具名 adopter 拉动,不抢跑",
                align=PP_ALIGN.LEFT)

    # 17 — 总结 / 闭幕
    s = add_blank_slide(prs)
    hdr(s, "一句话回顾", 17)
    add_filled_rect(s, Inches(0.8), Inches(2.2),
                    SLIDE_W - Inches(1.6), Inches(4.0),
                    fill=COLOR_PRIMARY)
    add_textbox(s, Inches(1.1), Inches(2.5),
                SLIDE_W - Inches(2.2), Inches(3.5),
                text=(
                    "ATD = 一份冻结的 5-message 中立协议\n"
                    "+ 一套可装配的 server runtime\n"
                    "  (capability · audit · rate · TokenBroker · UCAN-lite · Cursor · middleware)\n"
                    "+ 一组桥接 (MCP-bridge · SDK · CLI)\n\n"
                    "让 vendor 写一份 server,被任意 agent 平台用,\n"
                    "自带审计 / 多租户 / 跨 vendor 组合 / 子委托。"
                ),
                size=20, bold=False, color=COLOR_WHITE, line_spacing=1.45)
    add_caption(s, Inches(6.5),
                "下一份 deck: 01-design-philosophy — 7 条原则",
                align=PP_ALIGN.RIGHT,
                color=COLOR_MUTED)

    prs.save(str(out_path))
    return slides_total


def build_deck_01_philosophy(out_path: Path):
    deck_title = "01 · ATD 设计哲学"
    prs = new_presentation()
    add_cover_slide(
        prs,
        deck_title="ATD 设计哲学 ── 7 条原则 + 反模式 + Adopter Checklist",
        subtitle=(
            "Wire frame 是给 LLM 看的,audit sink 是给人看的,\n"
            "handshake 是给桥接用的。每个设计决策要同时通过三个读者检验。"
        ),
        footer="基于 docs/atd-design-philosophy.md · 2026-05-19",
    )

    slides_total = 21

    def hdr(s, title, n):
        add_page_header(s, deck_title=deck_title, slide_title=title,
                        page_num=n, total=slides_total)

    # 02 — 3 个消费者
    s = add_blank_slide(prs)
    hdr(s, "ATD tool server 同时面对 3 个消费者", 2)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(3.5),
        header=["消费者", "需求", "通道"],
        rows=[
            ["LLM Agent",
             "可发现的工具面 / 类型化错误 / 可预测的 arg shape",
             "tool_list · tool_schema · run_tool over wire"],
            ["人类运维",
             "审计轨迹 / 运维控制 / 结构化日志 / capability 拒绝可见性",
             "AuditSink 事件 · server log · metrics"],
            ["Agent 平台桥接",
             "稳定握手 / capability 协商 / 不出意外的传输",
             "Hello/HelloAck + UCAN-lite · length-prefixed JSON over UDS/HTTP/stdio"],
        ],
        cell_size=14,
        col_widths=[0.18, 0.42, 0.40],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.7),
                    SLIDE_W - Inches(1.0), Inches(1.4),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.85),
                SLIDE_W - Inches(1.4), Inches(1.15),
                text=(
                    "Wire frame = LLM。Audit sink = 人。Handshake = 桥接。\n"
                    "三条管子,同一个 server,没有 flag、没有 mode。\n"
                    "让 LLM 爽但让桥接握手崩 = bug,不是 trade-off。"
                ),
                size=15, italic=True, color=COLOR_PRIMARY, line_spacing=1.4)

    # 03 — 7 原则总览
    s = add_blank_slide(prs)
    hdr(s, "7 条原则 —— 一页总览", 3)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(5.4),
        header=["#", "原则", "一句话"],
        rows=[
            ["1", "ToolDefinition 是唯一真实源",
             "从一份 ToolDefinition 生成 summaries / 校验 / skills / adapter / 文档"],
            ["2", "Skill 跟着工具走,不跟着桥接走",
             "暴露 skills.list / skills.get;atd skills sync 按平台安装"],
            ["3", "Capability 协商而非硬编码",
             "声明 required;求交集 granted;dispatch 层 gate;handler 不查"],
            ["4", "Error 类型化、namespace 化",
             "协议 1000-1099; adopter 2000+;无自由文本主信号"],
            ["5", "工具默认跨连接无状态",
             "ConnectionContext per-conn; 共享世界 opt-in 且显式声明"],
            ["6", "Discovery canonical",
             "agent 运行时 discover; 不在 prompt 硬编码 tool id"],
            ["7", "Dispatch bounded + observable",
             "tier deadline; 中间件; 不静默重试; 失败可观测"],
        ],
        cell_size=13,
        col_widths=[0.05, 0.30, 0.65],
    )

    # 04 — 原则 1 详
    s = add_blank_slide(prs)
    hdr(s, "原则 1 — ToolDefinition 是唯一真实源", 4)
    left_body, right_body = add_two_panel(
        s, title_left="✅ 对的做法",
        title_right="❌ 反模式 — args 描述手维护",
    )
    add_code_block(s, *left_body, code=(
        "@server.register(\n"
        "  definition=ToolDefinition(\n"
        "    input_schema={\n"
        "      \"type\": \"object\",\n"
        "      \"properties\": {\n"
        "        \"path\": {\"type\": \"string\"}\n"
        "      },\n"
        "      \"required\": [\"path\"]\n"
        "    },\n"
        "    description=\"Read a file.\",\n"
        "  ))\n"
        "# Summary / SKILL.md / OpenAI shape\n"
        "# 全部从这份生成,无第二拷贝\n"
    ), size=12)
    add_code_block(s, *right_body, code=(
        "@server.register(\n"
        "  definition=ToolDefinition(\n"
        "    input_schema={...},\n"
        "    description=(\n"
        "      \"Read file from `path` \"\n"
        "      \"(required) and return \"\n"
        "      \"its contents.\"  # ← 会烂\n"
        "    ),\n"
        "  ))\n"
        "# 新增 optional encoding 时\n"
        "# 只有 schema 更新,\n"
        "# description 还说 \"from path\"\n"
    ), size=12)

    # 05 — 原则 2 详
    s = add_blank_slide(prs)
    hdr(s, "原则 2 — Skill 跟着工具走,不跟着桥接走", 5)
    add_textbox(s, Inches(0.5), Inches(1.85),
                SLIDE_W - Inches(1.0), Inches(0.6),
                text="SKILL.md 是工具的一部分,不是 agent 平台的一部分。",
                size=18, italic=True, color=COLOR_MUTED)
    add_code_block(s, Inches(0.5), Inches(2.6),
                   SLIDE_W - Inches(1.0), Inches(2.4),
                   code=(
                       "@server.register(definition=ToolDefinition(\n"
                       "    id=\"cbrain:sim.skills.list\",\n"
                       "    visibility=ToolVisibility.READ,\n"
                       "    required_capabilities=[],   # 公开 meta-tool\n"
                       "))\n"
                       "async def list_skills(args, ctx):\n"
                       "    return [{\"name\": p.parent.name,\n"
                       "             \"description\": _read_desc(p)}\n"
                       "            for p in sorted(SKILL_ROOT.glob(\"*/SKILL.md\"))]\n"
                   ),
                   size=13)
    add_code_block(s, Inches(0.5), Inches(5.2),
                   SLIDE_W - Inches(1.0), Inches(1.6),
                   code=(
                       "# 安装到 agent 平台\n"
                       "atd skills sync --target hermes\n"
                       "atd skills sync --target claude-code\n"
                       "# → ~/.hermes/skills/cbrain-sim-<name>/SKILL.md\n"
                   ),
                   size=13)

    # 06 — 原则 3 详
    s = add_blank_slide(prs)
    hdr(s, "原则 3 — Capability 协商而非硬编码", 6)
    add_textbox(s, Inches(0.5), Inches(1.85),
                SLIDE_W - Inches(1.0), Inches(0.6),
                text="硬编码 capability 在 handler 内部 → 4 个问题:",
                size=17, italic=True, color=COLOR_MUTED)
    add_bullet_list(s, Inches(0.5), Inches(2.4),
                    SLIDE_W - Inches(1.0), Inches(3.0),
                    items=[
                        "检查对 LLM 不可见 — tool_list 反映不出",
                        "不同 handler 漂移 (有的查 env, 有的查 header, 有的查 client_id)",
                        "Audit / observability 看到失败太晚 — handler 已经开始跑",
                        "未来想 pre-fetch capability 的桥接没东西可查 (如 UI 权限弹窗)",
                    ],
                    size=17)
    add_filled_rect(s, Inches(0.5), Inches(5.6),
                    SLIDE_W - Inches(1.0), Inches(1.4),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.75),
                SLIDE_W - Inches(1.4), Inches(1.15),
                text=(
                    "外化到 required_capabilities + dispatcher gate 解决全部四个:\n"
                    "LLM 在 tool_schema 看到要求 / 检查统一 / "
                    "audit 在 handler 跑之前看到 / 桥接可预取 schema"
                ),
                size=15, italic=True, color=COLOR_PRIMARY,
                line_spacing=1.4)

    # 07 — 原则 4 详
    s = add_blank_slide(prs)
    hdr(s, "原则 4 — Error 类型化、namespace 化", 7)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(3.6),
        header=["Code", "常量", "含义"],
        rows=[
            ["1000", "TOOL_NOT_FOUND", "工具不存在"],
            ["1001", "CAPABILITY_DENIED", "缺 required capability"],
            ["1002", "RATE_LIMITED", "per-tool semaphore 拒 (retryable)"],
            ["1003", "BROKER_FAILED", "TokenBroker 错"],
            ["1004", "DEADLINE_EXCEEDED", "超 tier deadline"],
            ["1005", "INVALID_ARGS", "args 不符 schema"],
            ["1010-1013", "UCAN_*", "UCAN invalid / expired / too-deep / aud-mismatch"],
            ["1020 / 1021", "CURSOR_*", "Cursor expired / invalid"],
            ["1099", "INTERNAL", "未捕获 (需人介入)"],
        ],
        cell_size=12,
        col_widths=[0.16, 0.32, 0.52],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.6),
                    SLIDE_W - Inches(1.0), Inches(1.4),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.75),
                SLIDE_W - Inches(1.4), Inches(1.15),
                text=(
                    "Adopter 占 2000+; cbrain 2000-2099 / healthkit 3000-3099 / celia 4000-4099\n"
                    "数字 code 跨翻译生存; 自由文本要 LLM 读散文,恢复不可靠"
                ),
                size=14, italic=True, color=COLOR_PRIMARY,
                line_spacing=1.4)

    # 08 — 原则 5 详
    s = add_blank_slide(prs)
    hdr(s, "原则 5 — 工具默认跨连接无状态", 8)
    left_body, right_body = add_two_panel(
        s, title_left="✅ Stateless (easy case)",
        title_right="✅ Shared world (declared)",
    )
    add_code_block(s, *left_body, code=(
        "@server.register(\n"
        "  definition=ToolDefinition(id=\"x:hash\"))\n"
        "async def hash(args, ctx):\n"
        "  return {\"sha256\":\n"
        "    hashlib.sha256(\n"
        "      args[\"data\"].encode()\n"
        "    ).hexdigest()}\n"
        "  # 不用 ctx, 不用模块状态\n"
    ), size=12)
    add_code_block(s, *right_body, code=(
        "SIM = MuJoCoSimulator()  # singleton\n\n"
        "@server.register(definition=ToolDefinition(\n"
        "  id=\"cbrain:manipulation.pick\",\n"
        "  description=(\n"
        "    \"Pick. WARNING: mutates \"\n"
        "    \"shared simulator state \"\n"
        "    \"visible to ALL agents.\"  # ← 显式\n"
        "  )))\n"
        "async def pick(args, ctx):\n"
        "  await SIM.pick(args[\"target\"])\n"
    ), size=12)

    # 09 — 原则 6 详
    s = add_blank_slide(prs)
    hdr(s, "原则 6 — Discovery canonical (不在 prompt 硬编码 tool id)", 9)
    add_filled_rect(s, Inches(0.5), Inches(1.85),
                    SLIDE_W - Inches(1.0), Inches(2.0),
                    fill=RGBColor(0xFE, 0xE2, 0xE2))
    add_textbox(s, Inches(0.7), Inches(1.95),
                SLIDE_W - Inches(1.4), Inches(1.8),
                text=(
                    "❌ System prompt:\n"
                    "\"You may call cbrain:perception.snapshot, cbrain:manipulation.pick,\n"
                    " cbrain:world.reset. Always start by calling perception.snapshot.\""
                ),
                size=15, color=RGBColor(0x7F, 0x1D, 0x1D),
                line_spacing=1.4)
    add_bullet_list(s, Inches(0.5), Inches(4.2),
                    SLIDE_W - Inches(1.0), Inches(2.8),
                    items=[
                        "cbrain-sim 加 cbrain:perception.depth_snapshot → prompt 不知道",
                        "manipulation.pick 改名 manipulation.grasp → 所有 agent 一起断",
                        "正确做法: agent session 启动时 discover() → 新工具自动出现, 改名不破坏",
                        "ToolSummary.id 是唯一稳定 handle, 其他都是人类面散文",
                    ],
                    size=16, line_spacing=1.4)

    # 10 — 原则 7 详
    s = add_blank_slide(prs)
    hdr(s, "原则 7 — Dispatch bounded + observable", 10)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(3.3),
        header=["contract", "实现", "失败模式"],
        rows=[
            ["Bounded",
             "definition.resources.timeout_ms (默认 30s)",
             "超时 → 1004 DEADLINE_EXCEEDED"],
            ["Observable",
             "middleware (pre/post/on_error)",
             "audit / tracing / metrics 均可挂"],
            ["No silent retries",
             "server 不内部重试",
             "transient → 返 retryable=true, client 决定"],
        ],
        cell_size=14,
        col_widths=[0.20, 0.40, 0.40],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.3),
                    SLIDE_W - Inches(1.0), Inches(1.7),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.45),
                SLIDE_W - Inches(1.4), Inches(1.4),
                text=(
                    "为什么 — 静默重试隐藏失败、对 side-effect 重复扣费;\n"
                    "不可观测 → 不知道是 cap-denied 还是 timed-out 还是 audit-dropped;\n"
                    "无 bound → agent 等桥接等 server 等工具,人页班"
                ),
                size=14, italic=True, color=COLOR_PRIMARY,
                line_spacing=1.4)

    # 11-17: 反模式速查 + Adopter checklist parts (split for visibility)
    # 11 — 反模式速查
    s = add_blank_slide(prs)
    hdr(s, "反模式速查 (10 条)", 11)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "手抄 SKILL.md 到 agent 平台 config 目录 (原则 2)",
                        "手写 args description 重复 input_schema (原则 1)",
                        "per-handler 硬编码 auth / capability 检查 (原则 3)",
                        "返回自由文本错误字符串不带数字 code (原则 4)",
                        "用 module-global state 模拟 per-connection (原则 5)",
                        "tool id 烧进 agent system prompt (原则 6)",
                        "Handler 内部隐式 retry loop (原则 7)",
                        "raise Exception(...) 作为主失败路径 (原则 4 + 7)",
                        "捕获 asyncio.CancelledError 后继续 (原则 7)",
                        "给 wire frame 加 per-platform shim (破坏 byte-compat)",
                    ],
                    size=16, line_spacing=1.4,
                    color=COLOR_BAD, bullet="✘  ")

    # 12 — Adopter Checklist 1 (schema + skills)
    s = add_blank_slide(prs)
    hdr(s, "Adopter Checklist [1/4] — schema + skills + discovery", 12)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "每个工具事实只在一份 ToolDefinition 里",
                        "ToolSummary 派生自 ToolDefinition, 不手维护",
                        "SKILL.md description frontmatter 与 ToolDefinition.description 一致",
                        "LLM-adapter shape (OpenAI / Anthropic) 从 summary 生成无平行映射",
                        "工具暴露 <publisher>:<service>.skills.list + .skills.get",
                        "agent 平台 config 内不手抄 SKILL.md (走 atd skills sync)",
                        "agent prompt 不硬编码 tool id (session 启动 discover)",
                    ],
                    size=17, line_spacing=1.5,
                    bullet="☐  ")

    # 13 — Adopter Checklist 2 (capabilities)
    s = add_blank_slide(prs)
    hdr(s, "Adopter Checklist [2/4] — capabilities", 13)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "每个需要 cap 的工具声明 required_capabilities",
                        "ServerPolicy 用 allow-list 求交集 (不 grant 一切)",
                        "Handler 内无 if not has_cap(...) 检查 — gating 在 dispatch",
                        "tool_schema 响应 include required_capabilities 给 LLM 看",
                    ],
                    size=17, line_spacing=1.5,
                    bullet="☐  ")

    # 14 — Adopter Checklist 3 (errors + state)
    s = add_blank_slide(prs)
    hdr(s, "Adopter Checklist [3/4] — errors + state", 14)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "每个 ToolError / ToolFailure 带数字 code (无 \"ERR\" / \"FAIL\")",
                        "Adopter code 落对 namespace (cbrain 2000+ / healthkit 3000+ / celia 4000+)",
                        "ToolDefinition.errors 广告该工具可能 emit 的 code",
                        "retryable 诚实 — True 只当 client 可安全 re-call",
                        "主失败路径无 raise Exception(...)",
                        "每个工具状态模型显式: stateless / per-connection / shared-world",
                        "Shared-world 工具在 description 说清楚",
                        "无 module-global 变量模拟 per-connection state",
                    ],
                    size=16, line_spacing=1.4,
                    bullet="☐  ")

    # 15 — Adopter Checklist 4 (observability + wire)
    s = add_blank_slide(prs)
    hdr(s, "Adopter Checklist [4/4] — observability + wire", 15)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "每个工具设 resources.timeout_ms (或自觉接 30s 默认)",
                        "中间件实现 audit / tracing / rate; dispatch 可观测",
                        "Handler 内无静默 retry loop (返 retryable=True)",
                        "asyncio.CancelledError 总是 re-raise",
                        "无 platform-specific shim 包 wire frame",
                        "跨语言实现 先过 atd-conformance fixture",
                    ],
                    size=17, line_spacing=1.5,
                    bullet="☐  ")

    # 16 — Adopter 检验示例
    s = add_blank_slide(prs)
    hdr(s, "原则在三家 adopter 的应用", 16)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.8),
        header=["原则", "healthkit_cli", "celia_phr", "cbrain"],
        rows=[
            ["1. ToolDefinition 源",
             "✅ HMS OpenAPI 单源", "✅ FHIR R4 单源", "🟡 SKILL home 风险"],
            ["2. Skills 跟工具走",
             "✅ skills.list/get + sync", "✅ 同上", "🟡 待迁移"],
            ["3. Capability 协商",
             "✅ records:read/write", "✅ 同模式", "🟡 待挑约定"],
            ["4. Error namespace",
             "✅ 3000-3099", "✅ 4000-4099", "🟡 待 claim 2000+"],
            ["5. Stateless",
             "✅ per-call token", "✅ HTTP per-request",
             "🟢 shared-world (intentional)"],
            ["6. Discovery canonical",
             "🟢 bridge-side discover", "🟢 同上", "🟢 同上"],
            ["7. Bounded + observable",
             "✅ timeout + retryable", "✅ audit mpsc + metrics",
             "🟢 Merkle middleware"],
        ],
        cell_size=12,
        col_widths=[0.22, 0.26, 0.26, 0.26],
    )

    # 17 — bigger picture
    s = add_blank_slide(prs)
    hdr(s, "ATD 是小协议 + 大约定", 17)
    add_textbox(s, Inches(0.5), Inches(1.95),
                SLIDE_W - Inches(1.0), Inches(2.5),
                text=(
                    "ATD 是一份小协议被一组大得多的设计选择包围。\n\n"
                    "协议 commit 给 wire 格式 / 握手 / discovery shape。\n"
                    "其他全是 adopter 约定 ——\n"
                    "如何发布 skill / 命名错误 / 限状态范围 / 设 deadline。"
                ),
                size=18, color=COLOR_TEXT, line_spacing=1.5)
    add_filled_rect(s, Inches(0.5), Inches(4.8),
                    SLIDE_W - Inches(1.0), Inches(2.3),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.0),
                SLIDE_W - Inches(1.4), Inches(2.0),
                text=(
                    "约定在软件里会静默退化。v1 schema-source-of-truth 在 v2 "
                    "被「再加一个 docstring」侵蚀。v3 ship 的 skills-via-meta-tools "
                    "在 v5 又静悄悄变回手抄。\n\n"
                    "约定不退化的方式: 写下来、用 adopter 示例佐证、"
                    "每次新 adopter 集成时重读。本文档就是这个 artifact。"
                ),
                size=14, italic=True, color=COLOR_PRIMARY,
                line_spacing=1.5)

    # 18 — 阅读顺序 (新 adopter)
    s = add_blank_slide(prs)
    hdr(s, "新 adopter 推荐阅读顺序", 18)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "1. docs/atd-positioning.md — ATD 是什么、为什么 (5 分钟)",
                        "2. docs/atd-architecture.md — 系统视图 (20 分钟)",
                        "3. docs/atd-design-philosophy.md — adopter 原则 (15 分钟)",
                        "4. docs/integrations/overview.md — 桥接 / 接入路径",
                        "5. docs/integrations/python-server.md 或 crates/atd-server/README.md — server runtime hello-world",
                        "6. 目标桥接的集成 recipe (hermes.md / claude-code.md / openclaw.md / ...)",
                        "7. docs/protocol/wire-format.md — 参考, 按需查",
                    ],
                    size=16, line_spacing=1.4)
    add_caption(s, Inches(6.9),
                "前 3 个是哲学,其余是执行。直接跳到执行就会出本文档要防的漂移。",
                align=PP_ALIGN.LEFT, italic=True)

    # 19 — 三个消费者的快速 reminder
    s = add_blank_slide(prs)
    hdr(s, "Reminder: 每个决策对三个读者都成立", 19)
    add_filled_rect(s, Inches(0.5), Inches(2.0),
                    Inches(4.0), Inches(4.5), fill=COLOR_BG_PANEL)
    add_textbox(s, Inches(0.7), Inches(2.1),
                Inches(3.6), Inches(0.5),
                text="LLM Agent", size=20, bold=True, color=COLOR_PRIMARY)
    add_bullet_list(s, Inches(0.7), Inches(2.7),
                    Inches(3.6), Inches(3.7),
                    items=["typed error", "predictable args",
                           "discovery 自动", "structured intent"],
                    size=13, line_spacing=1.5)

    add_filled_rect(s, Inches(4.7), Inches(2.0),
                    Inches(4.0), Inches(4.5), fill=COLOR_BG_PANEL)
    add_textbox(s, Inches(4.9), Inches(2.1),
                Inches(3.6), Inches(0.5),
                text="人类运维", size=20, bold=True, color=COLOR_PRIMARY)
    add_bullet_list(s, Inches(4.9), Inches(2.7),
                    Inches(3.6), Inches(3.7),
                    items=["audit trail (结构化 JSON Lines)",
                           "metrics counters",
                           "cap 拒绝可见",
                           "secret 不漏到 log"],
                    size=13, line_spacing=1.5)

    add_filled_rect(s, Inches(8.9), Inches(2.0),
                    Inches(4.0), Inches(4.5), fill=COLOR_BG_PANEL)
    add_textbox(s, Inches(9.1), Inches(2.1),
                Inches(3.6), Inches(0.5),
                text="桥接 (Hermes / MCP / 自研)",
                size=18, bold=True, color=COLOR_PRIMARY)
    add_bullet_list(s, Inches(9.1), Inches(2.7),
                    Inches(3.6), Inches(3.7),
                    items=["稳定握手", "cap 协商",
                           "transport 无意外",
                           "byte-compat wire frame"],
                    size=13, line_spacing=1.5)

    # 20 — 一句话回顾
    s = add_blank_slide(prs)
    hdr(s, "一句话回顾", 20)
    add_filled_rect(s, Inches(0.8), Inches(2.5),
                    SLIDE_W - Inches(1.6), Inches(3.5),
                    fill=COLOR_PRIMARY)
    add_textbox(s, Inches(1.1), Inches(2.85),
                SLIDE_W - Inches(2.2), Inches(2.8),
                text=(
                    "协议是 wire 格式 + 握手 + discovery shape。\n"
                    "其余全是 adopter 约定 —— 7 条原则保它不退化。\n\n"
                    "写下来 · 用 adopter 示例佐证 · 每次新接入时重读。"
                ),
                size=22, color=COLOR_WHITE, line_spacing=1.5,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_caption(s, Inches(6.5),
                "下一份 deck: 02-architecture-deepdive",
                align=PP_ALIGN.RIGHT)

    # 21 — references
    s = add_blank_slide(prs)
    hdr(s, "References", 21)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "docs/atd-design-philosophy.md (本 deck 主源,2026-05-19)",
                        "docs/atd-positioning.md (定位)",
                        "docs/atd-architecture.md (架构)",
                        "docs/protocol/error-codes.md (1000-1099 协议码完整表)",
                        "docs/issues/2026-05-19-cbrain-adopter-requirements.md (cbrain checklist 实证)",
                        "docs/issues/2026-05-26-atd-ts-sdk-adopter-requirements.md (atd-ts pre-design 应用 7 原则)",
                    ],
                    size=15)

    prs.save(str(out_path))
    return slides_total


def build_deck_02_architecture(out_path: Path):
    deck_title = "02 · ATD 架构深度"
    prs = new_presentation()
    add_cover_slide(
        prs,
        deck_title="ATD 架构深度 ── 统一 schema · dispatch · 安全 · middleware · 17 crate",
        subtitle=(
            "3 核心机制 (Schema / Dispatch / Security) + 2 扩展机制 (Bindings / Middleware)\n"
            "1.x wire 已冻结 · UDS 与 HTTP listener 共用同 dispatch entry"
        ),
        footer="基于 docs/atd-architecture.md · 2026-05-21",
    )

    slides_total = 24

    def hdr(s, title, n):
        add_page_header(s, deck_title=deck_title, slide_title=title,
                        page_num=n, total=slides_total)

    # 02 — what ATD is
    s = add_blank_slide(prs)
    hdr(s, "What ATD is — wire protocol + 类型化 RPC 表面", 2)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["维度", "现状的分裂", "ATD 的答案"],
        rows=[
            ["任意工具", "CLI / REST / MCP / native SDK 各 shape",
             "一份 ToolDefinition 映射多 binding"],
            ["任意平台", "Linux / macOS / iOS / Android / HMOS",
             "binding 选择在 server 侧 dispatch 时决定"],
            ["任意 agent", "Claude Code 吃不下 OpenAI shape",
             "所有 agent 同 SDK; adapter 渲 per-provider"],
            ["任意 framework", "LangChain ≠ MCP ≠ App Intent",
             "一份定义,多 framework consumer"],
        ],
        cell_size=14,
        col_widths=[0.20, 0.40, 0.40],
    )
    add_caption(s, Inches(6.6),
                "三个读者: 协议实现者 / 内部贡献者 / 决策者", italic=True)

    # 03 — unified schema
    s = add_blank_slide(prs)
    hdr(s, "§2 统一 schema — 唯一真实源", 3)
    add_textbox(s, Inches(0.5), Inches(1.95),
                SLIDE_W - Inches(1.0), Inches(0.7),
                text=(
                    "每个 wire 上的消息,每个方向,每个 transport(UDS / HTTP),"
                    "都序列化为 /atd-protocol-schema.json 定义的 shape。"
                ),
                size=15, italic=True, color=COLOR_MUTED, line_spacing=1.4)
    add_bullet_list(s, Inches(0.5), Inches(2.85),
                    SLIDE_W - Inches(1.0), Inches(4.3),
                    items=[
                        "Schema 从 atd-protocol Rust type 通过 schemars 生成",
                        "JSON Schema 2020-12 meta-schema 校验,CI gate drift",
                        "覆盖: Envelope / Handshake / Discovery / Invocation / "
                        "ToolDefinition 全字段 / 错误 taxonomy / 分页 / capability 协商",
                        "跨语言 SDK 自动 type-compatible (Rust / Python / 未来 TS / Go / Swift)",
                        "跨 transport parity: UDS 与 HTTP 共用同 dispatch_request entry point",
                        "1.0 schema 冻结为 1.x 稳定面: additive minor, breaking major",
                    ],
                    size=15, line_spacing=1.4)

    # 04 — sanitize
    s = add_blank_slide(prs)
    hdr(s, "Sanitization — tool id 跨 LLM 函数名 slot", 4)
    add_filled_rect(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(1.5),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(2.1),
                SLIDE_W - Inches(1.4), Inches(1.2),
                text=(
                    "tool id 含 ':' 和 '.' (如 ref:fs.read) → 破 LLM / MCP 函数名 slot\n\n"
                    "atd-sdk::sanitize 规范化双向映射: ref:fs.read ↔ ref_fs_read"
                ),
                size=16, color=COLOR_TEXT, line_spacing=1.5)
    add_textbox(s, Inches(0.5), Inches(3.8),
                SLIDE_W - Inches(1.0), Inches(0.6),
                text="两种形式同时出现在协议流量中 ——",
                size=15, italic=True, color=COLOR_MUTED)
    add_bullet_list(s, Inches(0.5), Inches(4.4),
                    SLIDE_W - Inches(1.0), Inches(2.6),
                    items=[
                        "canonical (ref:fs.read) — wire 上的形式,server / SDK / audit log",
                        "sanitised (ref_fs_read) — LLM tool-calling shape 内,OpenAI / Anthropic function 名",
                        "MCP bridge 应用同 sanitize 规则,tool id 不论落哪 slot 都明确",
                    ],
                    size=15, line_spacing=1.4)

    # 05 — 5 层模型 + 3 core / 2 ext
    s = add_blank_slide(prs)
    hdr(s, "§3 5 层 + 3 核心 + 2 扩展", 5)
    layers = [
        ("User intent", COLOR_MUTED),
        ("Agent framework (Claude / Hermes / LangChain / 自研)", COLOR_PRIMARY),
        ("Skills layer (adjacent)", COLOR_ACCENT),
        ("Client SDK", COLOR_PRIMARY),
        ("Dispatch (capability · tier · binding · cursor · middleware)",
         COLOR_ACCENT),
        ("Listener (UDS / HTTP+MCP)", COLOR_PRIMARY),
        ("Tool universe (bindings + 扩展点)", COLOR_MUTED),
    ]
    top = Inches(1.85)
    h = Inches(0.50)
    gap = Inches(0.13)
    for i, (label, color) in enumerate(layers):
        y = top + Emu(int((h + gap) * i))
        bg = COLOR_PRIMARY if i in (1, 3, 5) else (
            COLOR_ACCENT if i in (2, 4) else RGBColor(0xCB, 0xD5, 0xE1))
        fg = COLOR_WHITE if i in (1, 2, 3, 4, 5) else COLOR_TEXT
        add_filled_rect(s, Inches(0.5), y, SLIDE_W - Inches(1.0), h, fill=bg)
        add_textbox(s, Inches(0.7), y, SLIDE_W - Inches(1.4), h,
                    text=label, size=15, bold=True, color=fg,
                    anchor=MSO_ANCHOR.MIDDLE)
    add_caption(s, Inches(6.4),
                "3 core (Schema / Dispatch / Security) + 2 extension (Bindings / Middleware)",
                italic=True)

    # 06 — 6 wire variants
    s = add_blank_slide(prs)
    hdr(s, "§4.1 Wire — 6 个 request 变体", 6)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.6),
        header=["变体", "用途"],
        rows=[
            ["Hello",
             "握手 · client_id / requested_capabilities / 可选 ucan_tokens"],
            ["Ping", "心跳 · 服务器回 Pong"],
            ["ToolList",
             "Discovery · 返回 Vec<ToolSummary>, 按 DiscoverFilter 过滤"],
            ["ToolSchema", "单工具深 describe · 返回完整 ToolDefinition"],
            ["RunTool",
             "调用 · tool_id / args / CallOptions → ToolResultResponse"],
            ["RunToolContinue", "分页续传 · 带 opaque cursor"],
        ],
        cell_size=15,
        col_widths=[0.22, 0.78],
    )
    add_caption(s, Inches(6.7),
                "Wire = length-prefixed JSON over duplex byte stream",
                italic=True)

    # 07 — ToolDefinition
    s = add_blank_slide(prs)
    hdr(s, "§4.2 ToolDefinition — 工具完整声明", 7)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(5.4),
                   code=(
                       "pub struct ToolDefinition {\n"
                       "    pub id: String,                       // \"ref:fs.read\"\n"
                       "    pub name: String,\n"
                       "    pub description: String,\n"
                       "    pub version: String,\n"
                       "    pub capability: ToolCapability,       // domain · actions · intent_examples\n"
                       "    pub input_schema: serde_json::Value,  // JSON Schema 2020-12\n"
                       "    pub output_schema: serde_json::Value,\n"
                       "    pub bindings: Vec<ToolBinding>,       // Native / Cli / future\n"
                       "    pub safety: ToolSafety,               // Read/Write/Financial/Privacy/...\n"
                       "    pub resources: ToolResources,         // timeout_ms · max_concurrent\n"
                       "    pub trust: ToolTrust,                 // L0-L4 · signature\n"
                       "    pub visibility: ToolVisibility,       // Read/Write/Dangerous/Hidden\n"
                       "    pub required_capabilities: Vec<String>,\n"
                       "    pub tier: Option<ToolTier>,           // Hot / Warm / Cold\n"
                       "    pub errors: Vec<ToolErrorDef>,\n"
                       "}"
                   ),
                   size=13)

    # 08 — error taxonomy
    s = add_blank_slide(prs)
    hdr(s, "§4.3 Error taxonomy — 两层隔开", 8)
    add_textbox(s, Inches(0.5), Inches(1.95),
                SLIDE_W - Inches(1.0), Inches(2.0),
                text=(
                    "AtdError = client 侧 Rust enum (ToolNotFound / "
                    "InvalidArguments / CapabilityDenied / "
                    "BindingUnavailable / ...) — 无 numeric code\n\n"
                    "数字 wire code = ERR_* u16 常量在 atd_protocol::messages "
                    "— 落在 Response::Error.code"
                ),
                size=15, color=COLOR_TEXT, line_spacing=1.5)
    add_table(
        s, Inches(0.5), Inches(4.3),
        SLIDE_W - Inches(1.0), Inches(2.7),
        header=["Code 范围", "用途"],
        rows=[
            ["1000-1099", "协议级 (atd-protocol::messages)"],
            ["1010-1013", "UCAN (invalid / expired / too-deep / aud-mismatch)"],
            ["1020-1021", "Cursor (expired / invalid)"],
            ["2000+", "Adopter 区段 (cbrain 2000+ / healthkit 3000+ / celia 4000+)"],
        ],
        cell_size=13,
        col_widths=[0.25, 0.75],
    )

    # 09 — cursor pagination
    s = add_blank_slide(prs)
    hdr(s, "§4.4 Cursor 分页 — stateless HMAC 签名", 9)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "next_cursor: String — opaque, HMAC-SHA256 over CBOR payload",
                        "CursorPayload 绑定 (tool_id, caller_id, args_fingerprint, "
                        "page_index, issued_at_unix, server_session)",
                        "默认 TTL 5 分钟; wire cap 512 byte",
                        "不能跨 caller 重放 (caller_id 在签名内)",
                        "不能针对篡改 args 重放 (args_fingerprint 在签名内)",
                        "server 重启后 server_session 变 → 旧 cursor 失效 (1020 ERR_CURSOR_EXPIRED)",
                        "SDK call_all 自动走 chain · MergePolicy: ConcatArray / ConcatField / FirstPageOnly",
                        "工具 opt-in: 覆盖 Tool::supports_pagination + Tool::call_paginated",
                    ],
                    size=15, line_spacing=1.4)

    # 10 — dispatch pipeline
    s = add_blank_slide(prs)
    hdr(s, "§5 Dispatch pipeline — 确定 8 步", 10)
    steps = [
        "1. accept connection",
        "2. Hello 握手 (capability gate · 可选 UCAN verify)",
        "3. receive RunTool / RunToolContinue",
        "4. registry.get(tool_id)",
        "5. capability check (refuse if required ⊄ granted → 1001)",
        "6. tier-aware deadline + max_output_bytes 解析",
        "7. TokenBroker::resolve(caller_id) → CallContext::secrets",
        "8. binding.invoke(args, &ctx) [或 call_paginated when cursor]",
        "9. middleware pipeline (RedactPaths / FHIR / PII / ...)",
        "10. serialize ToolResultResponse + 可选 next_cursor",
    ]
    add_bullet_list(s, Inches(0.6), Inches(1.95),
                    SLIDE_W - Inches(1.2), Inches(5.0),
                    items=steps,
                    size=15, line_spacing=1.3, bullet="")
    add_caption(s, Inches(7.0),
                "UDS 和 HTTP listener 共用同 atd_runtime::dispatch::dispatch_request",
                italic=True)

    # 11 — SDK API
    s = add_blank_slide(prs)
    hdr(s, "§5.1 Core SDK API", 11)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["API", "用途", "SDK 形式"],
        rows=[
            ["discover", "枚举可见工具",
             "AtdClient::discover(filter) -> Vec<ToolSummary>"],
            ["describe", "取完整 ToolDefinition",
             "AtdClient::describe(tool_id) -> ToolDefinition"],
            ["call", "调用单 result",
             "AtdClient::call(tool_id, args, opts) -> ToolResult"],
            ["call_page", "单页",
             "AtdClient::call_page(id, args, Option<&cursor>, opts)"],
            ["call_all", "自动走 cursor chain",
             "AtdClient::call_all(id, args, CallAllOptions)"],
            ["ping / hello", "心跳 / 协商",
             "AtdClient::ping() / hello(client_id, caps)"],
        ],
        cell_size=12,
        col_widths=[0.14, 0.30, 0.56],
    )
    add_caption(s, Inches(6.65),
                "Python SDK at python/src/atd_client/ 镜像同 API (sync + async)",
                italic=True)

    # 12 — capability gate (2-fold)
    s = add_blank_slide(prs)
    hdr(s, "§5.2 Capability gate — 两机制组合", 12)
    left_body, right_body = add_two_panel(
        s, title_left="操作员 allow-list (字符串)",
        title_right="UCAN-lite token (additive)",
    )
    add_bullet_list(s, *left_body, items=[
        "server 启动声明 --grant-capability",
        "client Hello.requested_capabilities 求交集",
        "未 offer 的请求 silently dropped",
        "未请求的 offer 不 granted",
        "工具 required ⊄ granted → 1001",
    ], size=13, line_spacing=1.4)
    add_bullet_list(s, *right_body, items=[
        "Hello.ucan_tokens: Vec<String> (JWT compact)",
        "UCAN verifier 走 attenuation chain",
        "did:key + Ed25519 only",
        "granted = strings ∪ ucan_caps (联合不交集)",
        "UcanRevocationStore trait + 默认 5 链深度",
    ], size=13, line_spacing=1.4)

    # 13 — tier-aware deadline
    s = add_blank_slide(prs)
    hdr(s, "§5.3 Tier-aware deadline — latency / cost class 信号", 13)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(3.0),
        header=["Tier", "默认 deadline", "典型用途"],
        rows=[
            ["Hot", "sub-second", "同步无 side-effect 查询 (time, env)"],
            ["Warm", "seconds", "大多数工具 (file IO, shell, web fetch)"],
            ["Cold", "minutes", "慢导入 / 大导出 / 模型推理"],
        ],
        cell_size=15,
        col_widths=[0.15, 0.25, 0.60],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.2),
                    SLIDE_W - Inches(1.0), Inches(1.8),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.4),
                SLIDE_W - Inches(1.4), Inches(1.5),
                text=(
                    "Cursor-paginated 工具的 tier deadline per-page —— "
                    "Cold 工具可在长 wall-time 上流式输出而不违反 page 级 SLO\n\n"
                    "Per-call override: CallOptions::deadline_ms\n"
                    "Per-server override: --tier-override CLI flag"
                ),
                size=14, italic=True, color=COLOR_PRIMARY,
                line_spacing=1.4)

    # 14 — bindings
    s = add_blank_slide(prs)
    hdr(s, "§5.4 Bindings — 开放扩展点", 14)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(3.0),
        header=["Binding", "行为"],
        rows=[
            ["NativeBinding",
             "委托给同 Rust 进程的 Tool impl (每个内置工具默认)"],
            ["CliBinding",
             "派生子进程 · JSON args → argv · honor ctx.deadline 配 "
             "SIGTERM-then-SIGKILL"],
        ],
        cell_size=15,
        col_widths=[0.22, 0.78],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.3),
                    SLIDE_W - Inches(1.0), Inches(1.7),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.5),
                SLIDE_W - Inches(1.4), Inches(1.4),
                text=(
                    "Binding trait 开放: GrpcBinding / WasmBinding / McpBinding "
                    "都可实现同 Binding::invoke 签名。\n"
                    "v1 总路由到 ToolBinding 第一个; preferred_binding 是后续 dispatcher 小升级。"
                ),
                size=14, italic=True, color=COLOR_PRIMARY, line_spacing=1.4)

    # 15 — TokenBroker
    s = add_blank_slide(prs)
    hdr(s, "§5.5 TokenBroker — multi-tenant secret routing", 15)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(2.0),
                   code=(
                       "pub trait TokenBroker: Send + Sync {\n"
                       "    fn resolve(&self, caller_id: Option<&str>) -> ResolveFuture;\n"
                       "    fn resolve_bearer(&self, bearer: &str) -> ResolveBearerFuture;\n"
                       "    fn accepted_token_formats(&self) -> &'static [&'static str];\n"
                       "}"
                   ),
                   size=13)
    add_bullet_list(s, Inches(0.5), Inches(4.1),
                    SLIDE_W - Inches(1.0), Inches(3.0),
                    items=[
                        "InMemoryTokenBroker — 单元测试 / 单进程; UCAN-JWT 分支 register_ucan_audience()",
                        "FileTokenBroker — 磁盘后端; per-bearer subdir, 0700/0600; refresh mutex 防 OAuth 双 round-trip",
                        "HTTP bearer auth 走同 trait 的 resolve_bearer 臂 — 11 BearerOutcome 变体",
                        "SecretBundle wrap value 为 RedactedString — Debug/Display 拒打印",
                        "Audit event 只含 secrets_resolved: bool — 永不含 key 名或 value",
                    ],
                    size=14, line_spacing=1.4)

    # 16 — security 三轴
    s = add_blank_slide(prs)
    hdr(s, "§6.1 安全 — 三轴分类 (描述 metadata, 非执行机制)", 16)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(3.5),
        header=["分类", "值", "字段"],
        rows=[
            ["Safety",
             "Read / Write / Financial / Privacy / Physical / Destructive",
             "ToolSafety::level"],
            ["Visibility",
             "Read / Write / Dangerous / System / Hidden",
             "ToolVisibility"],
            ["Trust",
             "L0Unverified / L1SchemaValid / L2Tested / L3Verified / L4Certified",
             "ToolTrust::trust_level"],
        ],
        cell_size=13,
        col_widths=[0.15, 0.62, 0.23],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.6),
                    SLIDE_W - Inches(1.0), Inches(1.4),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.75),
                SLIDE_W - Inches(1.4), Inches(1.1),
                text=(
                    "Visibility::Hidden 把工具从 ToolList discovery 排除 "
                    "但保留 ToolSchema 和 RunTool 可达 — "
                    "用于 raw vendor endpoint / debug helper / 集成测试工具。"
                ),
                size=14, italic=True, color=COLOR_PRIMARY, line_spacing=1.4)

    # 17 — per-tool runtime control
    s = add_blank_slide(prs)
    hdr(s, "§6.3 Per-tool runtime 控制", 17)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.6),
        header=["控制", "适用", "位置"],
        rows=[
            ["SSRF guard (loopback/RFC1918/link-local/CGN/TEST-NET/0.0.0.0/8/v4-mapped; 重定向每跳重查)",
             "ref:web.fetch", "atd-tools-web/src/fetch.rs::check_ssrf"],
            ["Header allow-list (Accept/Accept-Lang/Referer/UA; Auth+Cookie 拒)",
             "ref:web.fetch", "atd-tools-web/src/fetch.rs::build_headers"],
            ["Must-read-before-edit (mtime+size 证明)",
             "ref:fs.edit", "atd-runtime/src/tracker.rs"],
            ["SIGTERM → grace → SIGKILL subprocess timeout",
             "ref:shell.exec / pwsh", "atd-tools-shell/src/shared.rs"],
            ["Per-tool semaphore (max_concurrent)",
             "全部", "atd-runtime/src/registry.rs"],
            ["Request-arg schema validation",
             "全部", "per-tool call impl + serde"],
        ],
        cell_size=11,
        col_widths=[0.40, 0.20, 0.40],
    )

    # 18 — audit
    s = add_blank_slide(prs)
    hdr(s, "§6.4 Audit — 结构化 CallEvent", 18)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(4.0),
                   code=(
                       "pub struct CallEvent {\n"
                       "    pub ts: String,                  // RFC3339\n"
                       "    pub call_id: String,\n"
                       "    pub tool_id: String,\n"
                       "    pub caller_id: Option<String>,\n"
                       "    pub granted_capabilities: Vec<String>,\n"
                       "    pub duration_ms: u64,\n"
                       "    pub outcome: Outcome,            // Success / ExecutionFailed / ...\n"
                       "    pub tier: String,\n"
                       "    pub dry_run: bool,\n"
                       "    pub schema_version: u32,         // 当前 2\n"
                       "    pub secrets_resolved: bool,      // 永不含 key 名/value\n"
                       "    pub cursor_page: Option<u32>,    // 1-based\n"
                       "}"
                   ),
                   size=13)
    add_caption(s, Inches(6.05),
                "JsonLinesAuditSink 用 bounded mpsc + 专用 drain task — on_call non-blocking; "
                "drops counter 通过 Server::metrics_snapshot 暴露",
                italic=True)

    # 19 — middleware pipeline
    s = add_blank_slide(prs)
    hdr(s, "§7 Middleware — 3 个内置 + 开放 trait", 19)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(3.4),
        header=["Middleware", "Crate", "作用"],
        rows=[
            ["RedactPathsMiddleware", "atd-runtime",
             "剥 / mask JSON Pointer 路径 (如 $HOME)"],
            ["FhirMiddleware", "atd-middleware-fhir",
             "75-URI 白名单 + 12-resource required field + 3 MismatchPolicy"],
            ["PiiRedactMiddleware", "atd-middleware-pii-redact-medical",
             "18 HIPAA × 13 JSON-Pointer × 7 RedactionStrategy + 5 catch-all 正则"],
        ],
        cell_size=12,
        col_widths=[0.30, 0.27, 0.43],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.5),
                    SLIDE_W - Inches(1.0), Inches(1.5),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.65),
                SLIDE_W - Inches(1.4), Inches(1.25),
                text=(
                    "I1 不变量: ALLOWED_SYSTEMS_DEFAULT 与 celia whitelists.toml set-equal\n"
                    "vendored toml + 双 CI gate; 任一仓单独改 set 就在两 gate 之一失败"
                ),
                size=14, italic=True, color=COLOR_PRIMARY, line_spacing=1.5)

    # 20 — 并发与限流
    s = add_blank_slide(prs)
    hdr(s, "§6.5 并发与限流", 20)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.0),
        header=["机制", "行为"],
        rows=[
            ["per-tool semaphore",
             "max_concurrent 在 Registry 强制 · 拒 RateLimited (1002 retryable)"],
            ["Multi-thread tokio",
             "ref binary 默认 multi_thread + min(cpus, 4) worker"],
            ["Per-state frame deadline (UDS)",
             "5s 握手 / 30s 活跃 · Server::set_frame_deadlines"],
            ["SDK connect retry",
             "指数退避 + ±20% jitter · ATD_CONNECT_RETRIES env"],
            ["Server 侧 rate-limiter",
             "v1 不在 · rate_limit_per_min declarative-only"],
        ],
        cell_size=13,
        col_widths=[0.30, 0.70],
    )
    add_caption(s, Inches(6.0),
                "50-client storm: p99=125ms 错=0 audit_drops=0 (vs pre-SP 71s wall + 60% 失败)",
                italic=True)

    # 21 — Skills layer
    s = add_blank_slide(prs)
    hdr(s, "§8 Skills layer (adjacent)", 21)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.3),
        header=["关注", "拥有者"],
        rows=[
            ["SKILL.md 创作 / 校验 / 安装", "Skills runtime (Anthropic / OpenClaw / 第三方)"],
            ["渐进披露到 agent context", "Skills runtime"],
            ["atd-tools: 依赖声明", "SKILL.md 格式; ATD 贡献是稳定 tool id"],
            ["从 skill body 调 ATD 工具", "Skills runtime 像任何 agent 调 SDK"],
            ["discover / describe / call API", "ATD (本仓)"],
        ],
        cell_size=13,
        col_widths=[0.40, 0.60],
    )
    add_caption(s, Inches(6.5),
                "ATD 不解析 SKILL.md, 不拥有平台安装路径, 不保留 skill state 跨调用",
                italic=True)

    # 22 — 17 crate map
    s = add_blank_slide(prs)
    hdr(s, "§9 17 crate workspace map", 22)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(5.4),
        header=["crate", "Layer", "职责"],
        rows=[
            ["atd-protocol", "Schema", "wire 格式 + 类型 + sanitize"],
            ["atd-sdk", "Client", "Rust 客户端 SDK"],
            ["atd-runtime", "Server core",
             "Tool · Registry · dispatch · Binding · Middleware · CursorIssuer · TokenBroker · UCAN · MetricsCounters"],
            ["atd-server", "Transport", "Unix socket listener"],
            ["atd-server-http", "Transport",
             "HTTP + MCP JSON-RPC translator + bearer + SSE"],
            ["atd-middleware-fhir", "Middleware",
             "FHIR R4 egress validation"],
            ["atd-middleware-pii-redact-medical", "Middleware",
             "HIPAA Safe Harbor PHI redaction"],
            ["atd-tools-{echo,fs,shell,web}", "Tools", "4 内置工具"],
            ["atd-mcp-bridge", "Bridge bin", "MCP/stdio → ATD"],
            ["atd-cli", "Bin", "atd 开发者 CLI"],
            ["atd-ref-server", "Bin", "参考 server binary"],
            ["atd-conformance", "Test", "跨实现 conformance"],
            ["atd-mock-weather-server", "Bin (publish=false)",
             "cross-vendor demo"],
        ],
        cell_size=11,
        col_widths=[0.32, 0.18, 0.50],
    )

    # 23 — 扩展点
    s = add_blank_slide(prs)
    hdr(s, "§9.3 扩展点 (不 fork ref-server 即可挂接)", 23)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.8),
        header=["想做…", "表面", "需 fork?"],
        rows=[
            ["加新工具", "Tool trait impl + Registry::register", "否"],
            ["加新 binding", "Binding trait impl", "否"],
            ["加新 middleware",
             "Middleware trait impl + Server::set_middleware", "否"],
            ["加新 auth scheme",
             "TokenBroker trait impl + ServerConfig::token_broker", "否"],
            ["加新 audit sink",
             "AuditSink trait impl + ServerConfig::audit_sink", "否"],
            ["加新 transport",
             "新 listener call atd_runtime::dispatch::dispatch_request",
             "否"],
            ["改 wire 格式", "—", "是 (不是扩展点)"],
        ],
        cell_size=13,
        col_widths=[0.20, 0.55, 0.25],
    )

    # 24 — non-goals
    s = add_blank_slide(prs)
    hdr(s, "§10 Non-goals (v1 故意不做)", 24)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "Multi-device routing (每连接 dispatch 到一个 socket)",
                        "Distributed session 迁移 / fork / handoff (scope 到一连接)",
                        "Tool signature verification (declarative; PKI 不规范)",
                        "REST / AppFunction / 分布式 binding (trait 可承载, 不 ship)",
                        "Native Skills-layer 支持 (Skills runtime 自己拥有)",
                        "Per-tool dry-run preview semantics (v1 server-side short-circuit)",
                        "Per-tool rate-limiter 强制 (declarative-only; 等 adopter 需要)",
                    ],
                    size=15, line_spacing=1.45, color=COLOR_MUTED,
                    bullet="✘  ")
    add_caption(s, Inches(6.85),
                "每个 non-goal 都有 rationale · 门槛是具体需求而非愿景",
                italic=True)

    prs.save(str(out_path))
    return slides_total


def build_deck_03_celia(out_path: Path):
    deck_title = "03 · celia_phr 案例深度"
    prs = new_presentation()
    add_cover_slide(
        prs,
        deck_title="celia_phr 案例 ── ATD 最复杂 adopter 的端到端解析",
        subtitle=(
            "本地优先 · 零知识 · 专利级 PHR · 3 shell × 4 binding × cross-vendor federation\n"
            "驱动了 5 个 ATD SP, 是 ATD 在生产规模下能站住的证明"
        ),
        footer="基于 celia_phr/docs/ARCHITECTURE.md + ATD 三份宪法 + 5 SP design",
    )

    slides_total = 22

    def hdr(s, title, n):
        add_page_header(s, deck_title=deck_title, slide_title=title,
                        page_num=n, total=slides_total)

    # 02 — celia 是什么
    s = add_blank_slide(prs)
    hdr(s, "celia_phr 是什么", 2)
    add_filled_rect(s, Inches(0.5), Inches(1.85),
                    SLIDE_W - Inches(1.0), Inches(1.7),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(2.0),
                SLIDE_W - Inches(1.4), Inches(1.4),
                text=(
                    "本地优先、零知识、专利级 Personal Health Record 应用。\n\n"
                    "DEK 只在易失内存 / 三 shell 一份 Rust core / "
                    "ATD 是其 agent 表面的中立调度协议"
                ),
                size=18, color=COLOR_TEXT, line_spacing=1.5)
    add_kv_grid(s, Inches(0.5), Inches(4.0),
                SLIDE_W - Inches(1.0), Inches(3.0),
                pairs=[
                    ("业务核心", "Rust ~41k LoC · 8 crate (auth/FHIR/crypto/vc/crdt/sub/RBAC)"),
                    ("Agent 表面", "celia binary — CLI + ATD UDS + atd-mcp-bridge + atd-server-http"),
                    ("三 shell 运行时", "Tauri 2.x 桌面 + Capacitor 6 移动 + PWA WebAssembly"),
                    ("存储", "SQLite (rusqlite native; PWA 走 sqlite-wasm-rs / IndexedDB)"),
                    ("专利", "§13.1 device-local DEK · §13.4 多 agent 隔离 · §13.5 multi-binding equivalence"),
                ],
                label_size=14, value_size=13)

    # 03 — 隐私不变量
    s = add_blank_slide(prs)
    hdr(s, "8 条隐私不变量 (每 PR review 抓)", 3)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(5.4),
        header=["不变量", "强制位置", "验证"],
        rows=[
            ["§13.1 DEK 只在易失内存",
             "celia-core/auth/key_cache.rs (Mutex<HashMap<UserId,Box<Zeroizing<[u8;32]>>>>)",
             "gcore 双 dump"],
            ["§13.1 Pattern A — DEK 父子只过 Unix socket",
             "src-tauri/agent_bootstrap.rs + celia-cli/parent_ipc.rs",
             "serve-pattern-a-test.sh"],
            ["§3 AES-256-GCM + SHA-256 双查",
             "fhir_store.rs::decrypt_and_verify", "171 cargo 测试"],
            ["§4 版本化 append + 软删",
             "FhirStore::create/update/soft_delete", "单元测试"],
            ["Coding 白名单 6 类",
             "celia-core/fhir/systems.rs (75-entry)",
             "写入 validation gate"],
            ["Agent gateway per-user in-process",
             "celia-tools::dispatch_for_caller", "3-dim RBAC"],
            ["Multi-agent isolation",
             "consent.grantee + CallerKind::External {agent_id}",
             "9 cargo 测试"],
            ["Multi-binding equivalence",
             "celia-tools transport-agnostic; 每 binding route 同 dispatch",
             "Phase J parity 6/6 + UDS↔HTTP byte-identical 2/2"],
        ],
        cell_size=11,
        col_widths=[0.30, 0.45, 0.25],
    )

    # 04 — 3 shell 单 core
    s = add_blank_slide(prs)
    hdr(s, "3 shell 单 Rust core 架构", 4)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(5.4),
                   code=(
                       "                 apps/web (React, 单 codebase)\n"
                       "                              │\n"
                       "                  services/celia-runtime.ts\n"
                       "                  ┌──────────┴──────────┐\n"
                       "                  │ isInTauri()         │\n"
                       "                  │ isInCapacitor()     │\n"
                       "                  │ isInBrowser()       │\n"
                       "                  └──────────┬──────────┘\n"
                       "       ┌─────────────────────┼─────────────────────┐\n"
                       "       ▼                     ▼                     ▼\n"
                       "  Tauri 2.x desktop    Capacitor 6 mobile     PWA / Browser\n"
                       "       │                     │                     │\n"
                       "  8 #[tauri::command]  UniFFI Swift/Kotlin    celia-core-wasm 1.4MB\n"
                       "       └─────────────────────┴─────────────────────┘\n"
                       "                              ▼\n"
                       " ╔══════════════════════════════════════════════════╗\n"
                       " ║  crates/celia-core (Rust 40k LoC, 单 source)     ║\n"
                       " ║  crypto + fhir + db + auth + vc + crdt + sub     ║\n"
                       " ╚══════════════════════════════════════════════════╝"
                   ),
                   size=10)

    # 05 — 21 工具 + 4 binding
    s = add_blank_slide(prs)
    hdr(s, "21 工具 · 4 binding 路径 · 同一份业务代码", 5)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(4.5),
                   code=(
                       "LLM agent\n"
                       "  │\n"
                       "  ├── (1) Tauri 命令 in-process   →  celia-tools::dispatch_for_caller\n"
                       "  ├── (2) ATD UDS                   →  atd_runtime::dispatch::run_tool\n"
                       "  ├── (3) MCP-stdio + atd-mcp-bridge → ATD UDS → 同上\n"
                       "  └── (4) MCP-over-HTTP + atd-server-http\n"
                       "                                    →  atd_runtime::dispatch::run_tool\n"
                       "\n"
                       "每条路径都通过相同的 RBAC + capability + audit + middleware\n"
                       "→ §13.1 / §13.4 / §13.5 在每条路径自动 hold"
                   ),
                   size=12)
    add_caption(s, Inches(6.5),
                "Parity test: tests/e2e_parity.rs 验证 UDS / HTTP 两 transport 对同一 call byte-identical",
                italic=True)

    # 06 — 4 SP 一览
    s = add_blank_slide(prs)
    hdr(s, "celia 触发的 5 个 ATD SP", 6)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.7),
        header=["SP", "celia 痛点", "ATD 答案"],
        rows=[
            ["SP-streamable-http",
             "云端 PHR 需 HTTP, UDS 跨 host 走不通",
             "atd-server-http · 同 dispatch entry"],
            ["SP-token-broker-phase2",
             "HTTP bearer · 多 token 形式 · SSE 长流",
             "resolve_bearer + 11 BearerOutcome + 60s SSE refresh"],
            ["SP-capability-v2 (UCAN-lite)",
             "Hermes orchestrator+N children · flat RBAC 强迫重 pair",
             "Hello.ucan_tokens · JWT Ed25519 did:key · attenuation chain"],
            ["SP-medical-middleware",
             "celia 已实现 FHIR + PHI 在错误层 · 应让其他 vendor 复用",
             "atd-middleware-fhir + atd-middleware-pii-redact-medical"],
            ["SP-concurrency-baseline",
             "10-concurrent benchmark 60% session-init failure",
             "multi_thread + frame deadline + SDK retry + audit mpsc"],
        ],
        cell_size=12,
        col_widths=[0.24, 0.40, 0.36],
    )

    # 07 — SP-streamable-http
    s = add_blank_slide(prs)
    hdr(s, "SP-streamable-http (1.B) — HTTP transport", 7)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.0),
                    items=[
                        "Trigger: celia 是云端可托管 PHR · Unix socket 跨 host 走不通",
                        "Ship: atd-server-http crate (axum + hyper)",
                        ("覆盖能力:", [
                            "MCP JSON-RPC translator",
                            "bearer auth (Authorization: Bearer ...)",
                            "origin gate (CSRF defense for browser MCP clients)",
                            "SSE refresh helper (60s 心跳 re-resolve token)",
                        ]),
                        "关键复用: 同一份 atd_runtime::dispatch::dispatch_request — UDS / HTTP 共享 dispatch 逻辑",
                        "v1.B 闭环: TLS termination / OAuth / 请求签名 留给 adopter (ATD 只管 transport + bearer 管线)",
                    ],
                    size=14, line_spacing=1.45)

    # 08 — SP-token-broker-phase2
    s = add_blank_slide(prs)
    hdr(s, "SP-token-broker-phase2 — multi-format bearer", 8)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["BearerOutcome 变体", "HTTP 映射"],
        rows=[
            ["Ok / OkShrunk", "200 + 正常 dispatch"],
            ["Expired", "401 + WWW-Authenticate error=invalid_token"],
            ["Revoked", "401 + ... error_description=revoked"],
            ["Unknown", "401 + ... error=invalid_token"],
            ["Internal", "500 + 不漏内部信息"],
            ["Lookup", "503 + Retry-After"],
            ["RateLimited / NetTimeout / Unsupported / ...",
             "对应 HTTP 状态码 + headers"],
        ],
        cell_size=13,
        col_widths=[0.40, 0.60],
    )
    add_caption(s, Inches(6.6),
                "SSE 60s 心跳 re-resolution → RefreshEvent::Refreshed / AuthLost",
                italic=True)

    # 09 — SP-capability-v2 overview
    s = add_blank_slide(prs)
    hdr(s, "SP-capability-v2 — UCAN-lite delegation", 9)
    add_textbox(s, Inches(0.5), Inches(1.95),
                SLIDE_W - Inches(1.0), Inches(2.0),
                text=(
                    "Trigger: Hermes orchestrator + N specialised children\n"
                    "  → 用户已经信任 Parent agent, 但 Parent 想让 child 看 patient X 3 个月只读\n"
                    "  → celia flat RBAC 强迫用户重 pair 每个 child (错的隐私姿态)\n\n"
                    "Keystone scenario: 「分享我近 3 个月心率给王医生, 7 天后失效」"
                ),
                size=15, color=COLOR_TEXT, line_spacing=1.5)
    add_bullet_list(s, Inches(0.5), Inches(4.6),
                    SLIDE_W - Inches(1.0), Inches(2.5),
                    items=[
                        "Hello.ucan_tokens: Vec<String> additive (pre-SP server 透明降级)",
                        "UCAN-lite v1.0 profile — JWT compact, Ed25519, did:key only",
                        "cmd=\"atd-cap\" 把 ATD 字符串 cap tunnel 进 UCAN payload",
                        "granted = strings ∪ ucan (联合不交集 — UCAN 已 attenuated)",
                        "max_ucan_chain_depth: u8 (默认 5, 可配; 防 verifier DoS)",
                    ],
                    size=14, line_spacing=1.4)

    # 10 — UCAN-lite payload 示例
    s = add_blank_slide(prs)
    hdr(s, "UCAN-lite payload 示例 — Agent A → child B", 10)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(4.5),
                   code=(
                       "{\n"
                       "  \"iss\": \"did:key:z<A's pubkey>\",\n"
                       "  \"aud\": \"did:key:z<B's pubkey>\",\n"
                       "  \"sub\": \"did:key:z<A's pubkey>\",\n"
                       "  \"cmd\": \"atd-cap\",\n"
                       "  \"args\": {\n"
                       "    \"caps\": [\"records:read\"],\n"
                       "    \"with\": [\n"
                       "      { \"patient\": \"Patient/abc123\" }\n"
                       "    ]\n"
                       "  },\n"
                       "  \"nonce\": \"<random 16 bytes, base64url>\",\n"
                       "  \"exp\": 1736208000\n"
                       "}"
                   ),
                   size=14)
    add_caption(s, Inches(6.5),
                "args.with 与 celia consent.patient_filter 1:1 — 同语义跨 UCAN / SQLite",
                italic=True)

    # 11 — UCAN-lite revocation
    s = add_blank_slide(prs)
    hdr(s, "UCAN-lite 撤销 — 2 层组合", 11)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.0),
        header=["层", "机制", "延迟"],
        rows=[
            ["Tier 1 (强制)", "TTL — 每 link 的 exp · 走得过期",
             "TTL 决定 worst-case (人发: 24h; 服务: ≤60s)"],
            ["Tier 2 (可选)",
             "UcanRevocationStore trait · 进程内 store · is_revoked(cid)",
             "Celia UI revoke → store.insert → 下次 Hello 失败 (ms 内)"],
        ],
        cell_size=13,
        col_widths=[0.18, 0.62, 0.20],
    )
    add_filled_rect(s, Inches(0.5), Inches(5.5),
                    SLIDE_W - Inches(1.0), Inches(1.5),
                    fill=COLOR_BG_TINT)
    add_textbox(s, Inches(0.7), Inches(5.65),
                SLIDE_W - Inches(1.4), Inches(1.25),
                text=(
                    "组合 SP-token-broker-phase2 §4.8 — "
                    "BrokerError::Revoked 复用同变体, 无新错误码。\n"
                    "SSE 长流: 60s 心跳 re-resolve, end-to-end 撤销最坏 ≤ 60s"
                ),
                size=13, italic=True, color=COLOR_PRIMARY, line_spacing=1.5)

    # 12 — UCAN-lite dormant
    s = add_blank_slide(prs)
    hdr(s, "UCAN-lite 当前状态 — shipped-dormant", 12)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(3.3),
        header=["维度", "状态"],
        rows=[
            ["代码", "✅ shipped end-to-end (27 unit + 12 integration green)"],
            ["生产流量", "🟡 仍走 ce_<hex> bearer · 因为无 adopter 真在 mint chain"],
            ["adopter 路径", "✅ celia 27 测试可启 (TokenBroker 不依赖 + 共存)"],
            ["唤醒触发", "🔮 当具名 adopter 真实启用 sub-agent delegation"],
        ],
        cell_size=14,
        col_widths=[0.20, 0.80],
    )
    add_caption(s, Inches(5.4),
                "Keystone scenario: 「分享我近 3 个月心率给王医生, 7 天后失效」 "
                "— 该场景产品化时即激活",
                italic=True)

    # 13 — SP-medical-middleware
    s = add_blank_slide(prs)
    hdr(s, "SP-medical-middleware — 2 个独立 crate", 13)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["Crate", "内容"],
        rows=[
            ["atd-middleware-fhir",
             "75-URI coding system 白名单 (与 celia source set-equal, drift-guard 双 CI gate) · "
             "12-resource Celia-subset required-field 校 · "
             "3 个 MismatchPolicy (AnnotateAndPass / ReplaceWithError / StripOffending)"],
            ["atd-middleware-pii-redact-medical",
             "18 HIPAA Safe Harbor identifier × 13 JSON-Pointer × "
             "7 RedactionStrategy (Strip/Token/FirstCharPrefix/HashSha256Truncated/YearOnly/ZipPrefix3) · "
             "5 catch-all 正则 (SSN/license/IP/URL/email) · "
             "fhir_aware opt-in"],
        ],
        cell_size=11,
        col_widths=[0.32, 0.68],
    )
    add_caption(s, Inches(6.6),
                "关键: CallEvent schema 不变 · PHI 永不出现在 audit · "
                "「audit 看到 redacted result」是 non-problem under v1",
                italic=True)

    # 14 — SP-concurrency-baseline 5 axis
    s = add_blank_slide(prs)
    hdr(s, "SP-concurrency-baseline — 5 轴干预 (全 back-compat)", 14)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["轴", "改"],
        rows=[
            ["Server runtime",
             "current_thread → multi_thread tokio · ATD_WORKER_THREADS env"],
            ["Wire deadline",
             "WireError::Timeout · read/write_frame_with_deadline · 5s 握手 / 30s 活跃"],
            ["SDK retry",
             "AtdClient::connect 指数退避 (5× · 50→800ms · ±20% jitter)"],
            ["Audit sink",
             "JsonLinesAuditSink 重写 bounded mpsc + 专用 drain · on_call non-blocking"],
            ["Metrics",
             "MetricsCounters + Server::metrics_snapshot() · 5 counter"],
        ],
        cell_size=13,
        col_widths=[0.20, 0.80],
    )
    add_caption(s, Inches(6.6),
                "celia adopter zero source edit · 仅 cargo rebuild 即生效",
                italic=True)

    # 15 — before / after
    s = add_blank_slide(prs)
    hdr(s, "SP-concurrency-baseline — before / after", 15)
    left_body, right_body = add_two_panel(
        s, title_left="❌ Pre-SP (incident)",
        title_right="✅ Post-SP",
    )
    add_bullet_list(s, *left_body, items=[
        "10-concurrent benchmark: 60% session-init failure",
        "6/10 sessions prompt_tokens ~180 (no-tools fallback)",
        "Hermes log: Connection lost ×N",
        "Wall clock ~71s",
    ], size=14, line_spacing=1.4)
    add_bullet_list(s, *right_body, items=[
        "Ref-server storm n=50: p99=125ms 错=0 audit_drops=0",
        "celia iter-4 SHARP baseline 120Q: 0 rate-limit / 0 connection 失败",
        "10/10 sessions full schema loaded (prompt_tokens ~5200)",
        "Wall 主由 DeepSeek LLM round-trip 决定 (非 ATD overhead)",
    ], size=14, line_spacing=1.4)

    # 16 — Phase L federation
    s = add_blank_slide(prs)
    hdr(s, "Phase L federation — celia 接 healthkit 当远程 ATD endpoint", 16)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(4.5),
                   code=(
                       "celia-connectors\n"
                       "  ├── AtdUpstreamIngest\n"
                       "  │     - atd-sdk + cursor 分页\n"
                       "  │     - Provenance (meta.source = atd://<endpoint>/<tool>)\n"
                       "  │\n"
                       "  ├── CursorStore trait\n"
                       "  │     - InMemory + FhirBasicCursorStore (CRDT max-by-advanced_at + Lamport)\n"
                       "  │\n"
                       "  └── SyncOrchestrator\n"
                       "        - tick scheduler + ±20% jitter\n"
                       "        - 指数退避 1m → 5m → 30m → 2h\n"
                       "        - 5-failure → Degraded\n"
                       "        - per-task tokio::spawn 失败隔离\n"
                       "        - audit events via atd-runtime::AuditSink"
                   ),
                   size=12)
    add_caption(s, Inches(6.5),
                "1020 ERR_CURSOR_EXPIRED → CursorStore::invalidate 标 tombstone, 下次从头拉",
                italic=True)

    # 17 — 端到端数据流 (part 1)
    s = add_blank_slide(prs)
    hdr(s, "端到端数据流 [1/2] — agent → dispatch", 17)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(5.4),
                   code=(
                       "1. 用户在 Tauri app 输入「给我看 patient X 最近 3 个月血压」\n"
                       "2. apps/web/services/agent-api.ts → invoke('celia_chat_stream', {...})\n"
                       "3. src-tauri/commands.rs → spawn Hermes orchestrator\n"
                       "4. Hermes orchestrator (LLM) sees 21 ATD tools + 19 SKILL.md\n"
                       "   pick 'celia:phr.observation.search'\n"
                       "   args = {patient_id: 'Patient/X', code: '85354-9 (BP)',\n"
                       "           date_range: '2026-02..2026-05'}\n"
                       "5. Hermes 发 RunTool over ATD UDS (Pattern A child socket)\n"
                       "6. atd_runtime::dispatch::dispatch_request\n"
                       "   ├── Hello cap gate: granted = ['records:read', 'patient:X']\n"
                       "   │   (UCAN-lite chain attenuated)\n"
                       "   ├── required = ['records:read']  ✓\n"
                       "   ├── tier = Warm, deadline = 5s\n"
                       "   ├── TokenBroker::resolve('hermes-orch')\n"
                       "   │     → SecretBundle{user_id, dek_ref}\n"
                       "   └── NativeBinding::invoke → celia_tools::dispatch_for_caller"
                   ),
                   size=11)

    # 18 — 端到端数据流 (part 2)
    s = add_blank_slide(prs)
    hdr(s, "端到端数据流 [2/2] — RBAC → middleware → audit", 18)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(5.4),
                   code=(
                       "7. celia-core RBAC check\n"
                       "   consent.grantee = 'agent:hermes-orch'\n"
                       "   consent.patient_filter = 'Patient/X' ✓\n"
                       "   consent.scope ⊇ ['records:read']    ✓\n"
                       "8. celia-core fhir_store\n"
                       "   SELECT ... WHERE patient_id='Patient/X' AND date BETWEEN ...\n"
                       "   for each row: decrypt_and_verify(encrypted_data, dek, hash)\n"
                       "9. Middleware pipeline (post-dispatch, on Value)\n"
                       "   ├── FhirMiddleware: resourceType ∈ {12} ✓ · system ∈ 75 ✓\n"
                       "   ├── PiiRedactMiddleware:\n"
                       "   │     /name → Token('NAME')\n"
                       "   │     /telecom → Token('PHONE')\n"
                       "   │     /address postalCode → ZipPrefix3\n"
                       "   │     /birthDate → YearOnly\n"
                       "   └── RedactPathsMiddleware (本例无 $HOME)\n"
                       "10. AuditSink::on_call\n"
                       "    CallEvent { ts, call_id, tool_id, caller_id, granted,\n"
                       "                duration_ms: 47, outcome: Success,\n"
                       "                tier: 'warm', secrets_resolved: true, ... }\n"
                       "    → JSONL via bounded mpsc → 不阻 dispatch"
                   ),
                   size=11)

    # 19 — 反事实分析
    s = add_blank_slide(prs)
    hdr(s, "反事实分析 — 不用 ATD celia 要自实现什么", 19)
    add_table(
        s, Inches(0.5), Inches(1.85),
        SLIDE_W - Inches(1.0), Inches(5.4),
        header=["自实现成本", "ATD 已 ship"],
        rows=[
            ["MCP / HTTP / Unix server 三套 transport 逻辑",
             "atd-server + atd-server-http + atd-mcp-bridge 共用 dispatch"],
            ["Capability gate 跨 transport 一致",
             "Hello.granted_capabilities + dispatch gate"],
            ["Audit log schema + 非阻塞写 + rotate",
             "JsonLinesAuditSink mpsc bounded"],
            ["Multi-tenant token routing",
             "TokenBroker::resolve + caller_id 路由"],
            ["OAuth bearer + SSE 心跳 re-validation",
             "resolve_bearer + sse_refresh"],
            ["UCAN-lite delegation (27 测 + 12 e2e)",
             "Hello.ucan_tokens + atd_runtime::ucan::*"],
            ["FHIR R4 validation + 75-URI 白名单",
             "atd-middleware-fhir"],
            ["HIPAA PHI redaction (18 类 × 13 路径)",
             "atd-middleware-pii-redact-medical"],
            ["Cursor 分页 + HMAC 签名 + cross-tool 重放防",
             "CursorIssuer + supports_pagination"],
            ["多 connection 并发 + p99 < 200ms SLO",
             "multi_thread tokio + frame deadline + SDK retry"],
            ["跨 transport byte-parity 测试",
             "atd-conformance"],
        ],
        cell_size=11,
        col_widths=[0.50, 0.50],
    )

    # 20 — celia 不要的
    s = add_blank_slide(prs)
    hdr(s, "celia 自己实现的 (ATD 故意 non-goal)", 20)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.7),
        header=["celia 自实现", "为何不在 ATD"],
        rows=[
            ["§13.1 DEK 加密 / KeyCache",
             "adopter-specific 加密策略; ATD operate on already-decrypted Value"],
            ["Pattern A IPC (Tauri 父子 socket bootstrap)",
             "桌面 app shape; ATD layer 不管 process model"],
            ["FHIR Bundle / CSV / JSON ingestion",
             "Ingestion 是产品 surface, 非协议"],
            ["CRDT 同步 + Lamport tiebreaker",
             "celia 业务模型"],
            ["WebCrypto PWA subset",
             "浏览器特定"],
            ["GDPR Article 17 erasure (hard delete)",
             "celia 业务模型 + 合规"],
            ["Tauri commands + UniFFI Swift/Kotlin",
             "桥接产品壳, 非协议"],
        ],
        cell_size=12,
        col_widths=[0.50, 0.50],
    )

    # 21 — celia 一句话总结
    s = add_blank_slide(prs)
    hdr(s, "celia 一句话总结", 21)
    add_filled_rect(s, Inches(0.5), Inches(2.0),
                    SLIDE_W - Inches(1.0), Inches(4.3),
                    fill=COLOR_PRIMARY)
    add_textbox(s, Inches(0.8), Inches(2.3),
                SLIDE_W - Inches(1.6), Inches(3.7),
                text=(
                    "celia_phr 证明:\n\n"
                    "一份协议级中立调度面 (ATD)\n"
                    "+ 一组可装配的扩展点 (Binding / Middleware / TokenBroker / AuditSink / UCAN)\n\n"
                    "能让一个本地优先、零知识、专利级 PHR 应用在\n"
                    "3 个 shell × 4 条 binding 路径 × cross-vendor federation × multi-agent delegation\n"
                    "场景下保持 single source of truth 的业务代码\n"
                    "+ 可逐路径验证的隐私不变量。"
                ),
                size=18, color=COLOR_WHITE, line_spacing=1.45)

    # 22 — references
    s = add_blank_slide(prs)
    hdr(s, "References", 22)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        "celia_phr/docs/ARCHITECTURE.md (2026-05-15 Phase L.4-followup-2)",
                        "docs/adr/0001-celia-atd-roadmap-alignment.md (4 family categorization)",
                        "docs/archive/superpowers/specs/2026-05-11-sp-streamable-http-design.md",
                        "docs/archive/superpowers/specs/2026-05-11-sp-token-broker-phase2-design.md",
                        "docs/archive/superpowers/specs/2026-05-11-sp-capability-v2-design.md",
                        "docs/archive/superpowers/specs/2026-05-11-sp-medical-middleware-design.md",
                        "docs/archive/superpowers/specs/2026-05-12-sp-concurrency-baseline-design.md",
                        "docs/issues/2026-05-12-celia-concurrency-adopter-validation.md (closed-verified)",
                        "docs/issues/2026-05-28-ucan-dormant-awaiting-product-trigger.md",
                    ],
                    size=14, line_spacing=1.4)

    prs.save(str(out_path))
    return slides_total


def build_deck_04_scenarios(out_path: Path):
    deck_title = "04 · ATD 高价值应用场景"
    prs = new_presentation()
    add_cover_slide(
        prs,
        deck_title="ATD 高价值应用场景 ── 7 种 raw 替代品做不到的事",
        subtitle=(
            "共通点: raw CLI / raw MCP / 自研 adapter 都能勉强能跑,\n"
            "但要生产质量都得重新实现 ATD ship 的中间层"
        ),
        footer="docs/intro/atd-tech-deck.zh.md §4",
    )

    slides_total = 13

    def hdr(s, title, n):
        add_page_header(s, deck_title=deck_title, slide_title=title,
                        page_num=n, total=slides_total)

    # 02 — 7 场景 overview
    s = add_blank_slide(prs)
    hdr(s, "7 个高价值场景一览", 2)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(5.0),
        header=["#", "场景", "典型 adopter"],
        rows=[
            ["1", "PHR / 医疗 vertical (合规 + 加密 + 多租户)",
             "celia_phr · healthkit_cli"],
            ["2", "Agent-Native CLI (一 binary 四出口)",
             "healthkit_cli · agentic-native-cli (mycli) · oh-cli"],
            ["3", "跨厂商工具组合",
             "cross-vendor-demo · healthkit + weather"],
            ["4", "Embodied agent / 物理仿真",
             "cbrain (MuJoCo + LLM)"],
            ["5", "跨设备 federation (远程 ATD endpoint)",
             "celia-connectors Phase L · healthkit upstream"],
            ["6", "Agentic IDE / 代码助手",
             "Cursor · Claude Code · 自研 IDE agent"],
            ["7", "多 agent 编排 (orchestrator + N children)",
             "Hermes 'Manager + Specialised Children'"],
        ],
        cell_size=12,
        col_widths=[0.05, 0.55, 0.40],
    )

    # 03 — 场景 1: PHR / 医疗
    s = add_blank_slide(prs)
    hdr(s, "场景 1 — PHR / 医疗 vertical", 3)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        ("FHIR R4 + HIPAA PHI 合规",
                         ["atd-middleware-fhir + atd-middleware-pii-redact-medical 即装即用",
                          "coding system whitelist 通过 drift-guard 跟 celia source set-equal"]),
                        ("多用户 / 多 OAuth token",
                         ["TokenBroker + caller_id 路由 · audit 只记 secrets_resolved: bool"]),
                        ("多 agent 子委托",
                         ["UCAN-lite Hello.ucan_tokens 走 attenuation chain · "
                          "Agent A → sub-agent B 「读 patient X 3 个月只读」"]),
                        ("跨厂商健康数据",
                         ["华为 HMS / Apple HealthKit / Garmin / Fitbit",
                          "每 vendor 一 ATD server, agent 看合并 catalog"]),
                        ("端到端加密不破",
                         ["DEK 留 Tauri 进程内存 · ATD 仅 dispatch 解密后 JSON · "
                          "middleware 在 egress 截 PHI"]),
                    ],
                    size=14, line_spacing=1.35)

    # 04 — 场景 2: Agent-Native CLI
    s = add_blank_slide(prs)
    hdr(s, "场景 2 — Agent-Native CLI (一 binary 四出口)", 4)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(2.6),
                   code=(
                       "mycli --command            # 1. 人类 CLI 子命令\n"
                       "mycli --atd-serve          # 2. ATD Unix server\n"
                       "atd-mcp-bridge → mycli     # 3. MCP-over-stdio (Claude / Cursor)\n"
                       "atd-server-http + mycli    # 4. MCP-over-HTTP (cloud / browser)"
                   ),
                   size=14)
    add_bullet_list(s, Inches(0.5), Inches(4.7),
                    SLIDE_W - Inches(1.0), Inches(2.3),
                    items=[
                        "给 LLM 暴露的不是 --help 散文而是结构化 description + intent_examples + input_schema",
                        "同一工具既给人也给 agent — 行为零分叉 (无 \"agent mode\" 二份代码)",
                        "skills.list / skills.get 让 SKILL.md 跟着 CLI 版本走",
                    ],
                    size=14, line_spacing=1.45)

    # 05 — 场景 3: cross-vendor
    s = add_blank_slide(prs)
    hdr(s, "场景 3 — 跨厂商工具组合", 5)
    add_code_block(s, Inches(0.5), Inches(1.85),
                   SLIDE_W - Inches(1.0), Inches(3.2),
                   code=(
                       "agent\n"
                       "  ├── socket A: /tmp/healthkit.sock\n"
                       "  │     huawei:hms.healthkit.heartrate.* (27 工具)\n"
                       "  │\n"
                       "  ├── socket B: /tmp/weather.sock\n"
                       "  │     openweather:atd-mock-weather.* (3 工具)\n"
                       "  │\n"
                       "  └── socket C: ... (任意数量)\n"
                       "\n"
                       "atd list → 合并 catalog · 30 工具"
                   ),
                   size=13)
    add_bullet_list(s, Inches(0.5), Inches(5.4),
                    SLIDE_W - Inches(1.0), Inches(1.6),
                    items=[
                        "CLI 不能在一进程里同时 +heartrate 和 +weather.now",
                        "每 vendor 自治: 自己的 audit / 自己的 broker / 自己的 cap allow-list",
                    ],
                    size=14)

    # 06 — 场景 4: Embodied agent
    s = add_blank_slide(prs)
    hdr(s, "场景 4 — Embodied agent / 物理仿真 (cbrain)", 6)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        ("为何 ATD Python server 是对的",
                         ["MjData singleton + LLM agent 必须 in-process — Rust server 跨 socket 太慢",
                          "多 agent 看同一份物理 (shared world)，通过 description 显式声明"]),
                        ("Merkle audit 通过 middleware",
                         ["每次 pick / place / reset 都进 chain · post_call 是观察点"]),
                        ("ATD 价值 — cbrain 触发了 SP-server-py-v1",
                         ["Python 进程的 atd_server 与 Rust runtime byte-compat (22/24 conformance)",
                          "同一份 wire 协议 · 物理仿真也能被任何 ATD client 调"]),
                        ("Bounded: resources.timeout_ms = 2000",
                         ["仿真步 ≤ 2s · LLM 不会卡住"]),
                    ],
                    size=14, line_spacing=1.4)

    # 07 — 场景 5: Federation
    s = add_blank_slide(prs)
    hdr(s, "场景 5 — 跨设备 federation (celia 接 healthkit)", 7)
    add_bullet_list(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    items=[
                        ("AtdUpstreamIngest 通过 atd-sdk + cursor 分页拉远程 FHIR",
                         ["每条记录加 meta.source = atd://<endpoint>/<tool> provenance"]),
                        ("CursorStore — 持久化为 FHIR Basic 资源",
                         ["CRDT max-by-advanced_at + Lamport tiebreaker"]),
                        ("SyncOrchestrator — tick scheduler",
                         ["±20% jitter · 指数退避 1m→5m→30m→2h · 5-failure → Degraded",
                          "per-task tokio::spawn 失败隔离"]),
                        ("远程 server 重启返回 1020 ERR_CURSOR_EXPIRED",
                         ["CursorStore::invalidate 标 tombstone · 重启拉"]),
                        ("Vendor 中性",
                         ["celia-connectors 不知道 healthkit · 任何 ATD-speaking server 都能接入"]),
                    ],
                    size=13, line_spacing=1.4)

    # 08 — 场景 6: Agentic IDE
    s = add_blank_slide(prs)
    hdr(s, "场景 6 — Agentic IDE / 代码助手", 8)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["痛点", "ATD 答案"],
        rows=[
            ["每个 IDE 重复实现 fs / shell / web 工具",
             "ref:fs.* / ref:shell.* / ref:web.* 一份 server 喂多 IDE"],
            ["不知道 agent 跑了哪些 shell 命令",
             "AuditSink JSON Lines · 团队可 review"],
            ["agent 默认可写, 误删风险",
             "Capability 默认只读 · 需要时升权"],
            ["建议 vs 执行混在一起",
             "dry-run 短路 · 拿 preview 决定真跑"],
        ],
        cell_size=14,
        col_widths=[0.40, 0.60],
    )

    # 09 — 场景 7: 多 agent 编排
    s = add_blank_slide(prs)
    hdr(s, "场景 7 — 多 agent 编排 (orchestrator + N children)", 9)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["痛点", "ATD 答案"],
        rows=[
            ["Parent 委托 child 不给 child 全部权限",
             "UCAN-lite delegation chain · attenuation 自动收缩"],
            ["Child 失败不让 parent 的 audit 混乱",
             "caller_id 隔离 audit"],
            ["Child 共享 parent 的 OAuth token",
             "TokenBroker 按 caller_id 路由"],
            ["不同 child 走不同 vendor",
             "Cross-vendor 合并 catalog"],
        ],
        cell_size=14,
        col_widths=[0.45, 0.55],
    )
    add_caption(s, Inches(6.6),
                "Keystone: 「分享我近 3 个月心率给王医生, 7 天后失效」 — UCAN-lite 的产品化触发",
                italic=True)

    # 10 — adopter 总览
    s = add_blank_slide(prs)
    hdr(s, "Active adopters 当前一览 (post-1.0)", 10)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["Adopter", "覆盖场景", "状态"],
        rows=[
            ["healthkit_cli", "1 + 2 + 3 + 5", "case study v1.4.0 头对头"],
            ["celia_phr", "1 + 5 + 7", "closed-verified"],
            ["cbrain", "4", "adopter-validation cycle"],
            ["oh-cli / mycli (ANC)", "2 + 6", "mcp-native-v1 落地"],
        ],
        cell_size=13,
        col_widths=[0.25, 0.30, 0.45],
    )

    # 11 — 三个 dormant
    s = add_blank_slide(prs)
    hdr(s, "Dormant — 等具名 adopter trigger 才激活的能力", 11)
    add_table(
        s, Inches(0.5), Inches(1.95),
        SLIDE_W - Inches(1.0), Inches(4.5),
        header=["能力", "状态", "触发条件"],
        rows=[
            ["UCAN-lite delegation",
             "shipped end-to-end · 生产仍走 bearer",
             "具名 adopter 真启 sub-agent delegation"],
            ["atd-ts SDK",
             "pre-design research 1107 行落",
             "具名 ArkTS agent-runtime adopter 出现"],
            ["SP-cli-dispatcher-v1",
             "wire CliBindingConfig 1.1.0 已 ship 类型",
             "具名 declarative-manifest adopter"],
        ],
        cell_size=13,
        col_widths=[0.25, 0.35, 0.40],
    )
    add_caption(s, Inches(6.6),
                "ATD post-1.0 SP 触发模式: 具名 adopter 拉动, 不抢跑",
                italic=True)

    # 12 — 共通点 — ATD ship 的中间层
    s = add_blank_slide(prs)
    hdr(s, "7 场景的共通点 — ATD 一次性 ship 的中间层", 12)
    add_filled_rect(s, Inches(0.5), Inches(1.95),
                    SLIDE_W - Inches(1.0), Inches(5.2),
                    fill=COLOR_BG_TINT)
    add_bullet_list(s, Inches(0.7), Inches(2.1),
                    SLIDE_W - Inches(1.4), Inches(4.8),
                    items=[
                        "Capability gate · 跨 transport 一致 · dispatch 层",
                        "Audit log · 结构化 JSON Lines · 非阻塞 mpsc · metrics 暴露",
                        "Multi-tenant token routing · caller_id · OAuth bearer · UCAN-lite",
                        "Cursor 分页 · HMAC 签名 · cross-tool 重放防",
                        "Middleware pipeline · 中间件 trait 开放 (FHIR / PHI 等可 mix-and-match)",
                        "Bindings trait 开放 (Native / Cli / 未来 Grpc / Wasm / Mcp)",
                        "Skill 同步 · skills.list / skills.get + atd skills sync",
                        "Cross-vendor 桥接 · 多 socket · 合并 catalog",
                    ],
                    size=15, line_spacing=1.45, color=COLOR_PRIMARY)

    # 13 — 闭幕
    s = add_blank_slide(prs)
    hdr(s, "ATD 的本质 — 写一次,被任意 agent 平台用", 13)
    add_filled_rect(s, Inches(0.5), Inches(2.0),
                    SLIDE_W - Inches(1.0), Inches(4.3),
                    fill=COLOR_PRIMARY)
    add_textbox(s, Inches(0.8), Inches(2.4),
                SLIDE_W - Inches(1.6), Inches(3.5),
                text=(
                    "ATD = 一份冻结的 5-message 中立协议\n"
                    "+ 一套可装配的 server runtime (8 中间层)\n"
                    "+ 一组桥接 (MCP-bridge · SDK · CLI)\n\n"
                    "vendor 写一份 server, 被任意 agent 平台用,\n"
                    "自带审计 / 多租户 / 跨 vendor 组合 / 子委托。\n\n"
                    "—— 这是协议层差异, 不是工具能力差异。"
                ),
                size=20, color=COLOR_WHITE, line_spacing=1.5,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    prs.save(str(out_path))
    return slides_total


# ─────────────────────────────── main ───────────────────────────────────


def main():
    out_dir = Path(__file__).parent
    out_dir.mkdir(parents=True, exist_ok=True)

    builders = [
        ("00-atd-overview.zh.pptx", build_deck_00_overview),
        ("01-design-philosophy.zh.pptx", build_deck_01_philosophy),
        ("02-architecture-deepdive.zh.pptx", build_deck_02_architecture),
        ("03-celia-phr-case-study.zh.pptx", build_deck_03_celia),
        ("04-scenarios.zh.pptx", build_deck_04_scenarios),
    ]

    results = []
    for name, fn in builders:
        out = out_dir / name
        count = fn(out)
        size_kb = out.stat().st_size // 1024
        results.append((name, count, size_kb))
        print(f"  ✓ {name:40s}  {count:3d} slides  {size_kb:4d} KB")

    print()
    print(f"Generated {len(results)} decks · "
          f"{sum(r[1] for r in results)} slides total · "
          f"{sum(r[2] for r in results)} KB total")


if __name__ == "__main__":
    main()
